import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog

from src.ui.message_dialog import ask_yes_no, show_warning
from typing import Dict, List, Optional, Tuple
from PIL import Image, ImageTk
import io
import os

from src.models.config import Settings
from src.ui.toolbar import Toolbar
from src.ui.status_bar import StatusBar
from src.ui.theme import UITheme
from src.ui.chrome import SectionHeader
from src.utils.annotation_text import format_annotation_list_preview


class AnnotationMarker:
    """批注标记（小方块 + 双击弹出内容）"""

    MARKER_SIZE = 22
    POPUP_OFFSET_X = 28
    POPUP_WIDTH = 300
    POPUP_MAX_HEIGHT = 420
    TITLE_AREA = 42
    POPUP_PADDING = 10

    def __init__(
        self,
        x: int,
        y: int,
        text: str,
        color: str = "#7C3AED",
        *,
        display_mode: str = "marker",
        original_text: str = "",
        placement: str = "right",
        box_width: int = 0,
        box_height: int = 0,
        source_x: int = None,
        source_y: int = None,
        text_orientation: str = "horizontal",
        font_size: int = 12,
        font_family: str = "",
        style_kind: str = "inline",
    ):
        self.x = x
        self.y = y
        self.text = text
        self.color = color
        self.display_mode = display_mode
        self.original_text = original_text
        self.placement = placement
        self.box_width = box_width
        self.box_height = box_height
        self.source_x = source_x
        self.source_y = source_y
        self.text_orientation = text_orientation or "horizontal"
        self.font_size = font_size
        self.font_family = font_family or "Microsoft YaHei"
        self.style_kind = style_kind

        self.collapsed_width = self.MARKER_SIZE
        self.collapsed_height = self.MARKER_SIZE
        self.expanded_width = self.POPUP_WIDTH

        self.is_expanded = False
        self.width = self.collapsed_width
        self.height = self.collapsed_height
        self.popup_x = 0
        self.popup_y = 0

        self.canvas_id = None
        self.text_id = None
        self.drag_data = {"x": 0, "y": 0}
        self.icon_rect = (0, 0, 0, 0)
        self.popup_rect = None
        self.close_rect = None
        self.popup_window_id = None

    @staticmethod
    def _wrap_font():
        import tkinter.font as tkfont

        return tkfont.Font(family="Microsoft YaHei UI", size=11)

    def _measure_body_height(self) -> int:
        """按与弹层 Textbox 相同的字体/宽度测算正文高度"""
        font = self._wrap_font()
        wrap_width = self.POPUP_WIDTH - 28
        line_height = font.metrics("linespace") + 6
        lines = 0

        for paragraph in self.text.split("\n"):
            if not paragraph.strip():
                lines += 1
                continue
            line = ""
            for ch in paragraph:
                trial = line + ch
                if font.measure(trial) <= wrap_width:
                    line = trial
                else:
                    if line:
                        lines += 1
                    line = ch
            if line:
                lines += 1

        # 留一点余量，避免测算偏矮导致画出白框
        return max(lines, 1) * line_height + 16

    def _calc_popup_height(self) -> int:
        """弹层总高度 = 标题区 + 正文区（超出最大高度时正文区可滚动）"""
        max_body = self.POPUP_MAX_HEIGHT - self.TITLE_AREA - self.POPUP_PADDING
        body_h = min(self._measure_body_height(), max_body)
        total = self.TITLE_AREA + body_h + self.POPUP_PADDING
        return int(max(min(total, self.POPUP_MAX_HEIGHT), 100))

    def toggle_expand(self):
        """切换展开/折叠状态"""
        self.is_expanded = not self.is_expanded
        if self.is_expanded:
            self.height = self._calc_popup_height()
        else:
            self.width = self.collapsed_width
            self.height = self.collapsed_height
            self.popup_rect = None
            self.close_rect = None

    def collapse(self):
        """折叠"""
        self.is_expanded = False
        self.width = self.collapsed_width
        self.height = self.collapsed_height
        self.popup_rect = None
        self.close_rect = None

    def refresh_expanded_size(self):
        """文本更新后刷新展开尺寸"""
        if self.is_expanded:
            self.height = self._calc_popup_height()


class App(ctk.CTk):
    """主应用窗口"""

    PPT_IMPORT_ZOOM = 0.6

    def __init__(self, settings: Settings):
        super().__init__()

        self.settings = settings
        self.selected_files: List[str] = []
        self.current_file_index: int = -1

        # PDF渲染相关
        self.pdf_doc = None
        self.current_page = 0
        self.total_pages = 0
        self.zoom_level = 1.0
        self.page_image = None
        self.tk_image = None
        self.text_positions = {}  # LiteParse / PPTX 文本位置信息
        self.ocr_text_positions: Dict[int, List[dict]] = {}  # 页内文字块缓存（视觉或 OCR）
        self.ocr_text_sources: Dict[int, str] = {}  # 缓存来源：vision / ocr
        self.ppt_slide_emu: Tuple[int, int] = (9144000, 6858000)
        self._preview_shell_ready = False
        self._page_image_id = None
        self._page_offset_x = 0
        self._page_offset_y = 0
        self._canvas_configure_after = None

        # 批注相关
        self.annotations: Dict[int, List[AnnotationMarker]] = {}  # 当前文件的批注
        self.annotations_by_file: Dict[str, Dict[int, List[AnnotationMarker]]] = {}
        self.preview_ink_by_file: Dict[str, Dict[int, List[dict]]] = {}
        self.project_file_path: str = None
        self.converted_pdf_paths: Dict[str, str] = {}
        self._autosave_job = None
        self.selected_marker: AnnotationMarker = None
        self.dragging = False
        self._press_pos: Tuple[float, float] = None
        self._did_drag = False
        self._inline_edit_marker: AnnotationMarker = None
        self._inline_edit_toplevel: ctk.CTkToplevel = None
        self._inline_edit_widget = None
        self._inline_edit_host = None
        self._inline_edit_opening = False
        self._inline_editor_width = 0
        self._inline_editor_height = 0
        self._undo_restoring = False
        self._drag_undo_snap = None
        from src.utils.undo_stack import UndoStack

        self._undo_stack: UndoStack = UndoStack(max_size=50)

        # 画布绘图（与 Web 预览墨迹共用 preview_ink_by_file）
        self.canvas_tool = "select"
        self.ink_color = "#ef4444"
        self._current_ink_stroke: Optional[dict] = None
        self._ink_tool_buttons: Dict[str, ctk.CTkButton] = {}

        # 配置窗口
        self.title("TO PDF · 中文批注")
        self.geometry("1440x920")
        self.minsize(1100, 720)

        UITheme.install()
        UITheme.apply_root(self)

        from src.utils.branding import apply_window_icon

        apply_window_icon(self)

        # 创建 UI 组件
        self._create_widgets()
        self._apply_theme()

        # 配置布局
        self._configure_layout()

        # 绑定快捷键
        self.bind("<Control-plus>", lambda e: self._zoom_in())
        self.bind("<Control-minus>", lambda e: self._zoom_out())
        self.bind("<Control-0>", lambda e: self._zoom_reset())
        self.bind_all("<BackSpace>", self._on_backspace_delete_annotation)
        self.bind_all("<Delete>", self._on_backspace_delete_annotation)
        self.bind_all("<Control-z>", self._on_undo)
        self.bind_all("<Control-Z>", self._on_undo)

        from src.services.web_preview_server import WebPreviewServer

        self.web_preview = WebPreviewServer()
        self.web_preview.start()
        self._web_preview_opened = False

        self.protocol("WM_DELETE_WINDOW", self._on_close)
        self.after(200, self._restore_last_session)

    def _create_widgets(self) -> None:
        """创建 UI 组件"""
        # 工具栏
        self.toolbar = Toolbar(self)

        # 主内容区域
        self.content_frame = ctk.CTkFrame(self, fg_color="transparent")

        # 左侧文件列表
        self.file_list_frame = ctk.CTkFrame(self.content_frame, width=248)
        self._create_file_list()

        # 中间预览区域
        self.preview_frame = ctk.CTkFrame(self.content_frame)

        # 右侧批注面板
        self.sidebar_frame = ctk.CTkFrame(self.content_frame, width=380)
        self._create_sidebar()

        # 状态栏
        self.status_bar = StatusBar(self)

    def _create_file_list(self) -> None:
        """创建文件列表"""
        header = SectionHeader(self.file_list_frame, text="文件列表")
        header.pack(fill="x", padx=UITheme.PAD, pady=(UITheme.PAD, UITheme.PAD_SM))

        # 文件列表
        self.file_listbox = ctk.CTkScrollableFrame(
            self.file_list_frame,
            fg_color=UITheme.FILE_LIST_BG,
            label_fg_color=UITheme.FILE_LIST_BG,
            **UITheme.scrollable_frame_kwargs(),
        )
        self.file_listbox.pack(fill="both", expand=True, padx=UITheme.PAD, pady=UITheme.PAD_SM)
        UITheme.style_scrollable_frame(self.file_listbox)

        self.file_hint = ctk.CTkLabel(self.file_list_frame, text="点击顶部「导入」添加 PDF / PPT")
        self.file_hint.pack(pady=UITheme.PAD)

    def _create_sidebar(self) -> None:
        """创建批注侧边栏（单卡片统一宽度，避免列表与编辑区不齐）"""
        pad = UITheme.PAD_SM
        self._editor_card = ctk.CTkFrame(self.sidebar_frame)
        self._editor_card.pack(fill="both", expand=True, padx=UITheme.PAD, pady=UITheme.PAD)
        UITheme.style_card(self._editor_card, elevated=True)

        title_row = ctk.CTkFrame(self._editor_card, fg_color="transparent")
        title_row.pack(fill="x", padx=pad, pady=(pad, 4))

        SectionHeader(title_row, text="批注管理").pack(side="left")

        from src.models.annotation_preset import get_preset_labels

        header_actions = ctk.CTkFrame(title_row, fg_color="transparent")
        header_actions.pack(side="right")

        delete_all_btn = ctk.CTkButton(
            header_actions,
            text="删除全部",
            width=72,
            height=28,
            command=self._delete_all_page_annotations,
        )
        delete_all_btn.pack(side="left", padx=(0, 6))
        UITheme.style_soft_danger(delete_all_btn)

        add_row = ctk.CTkFrame(header_actions, fg_color="transparent")
        add_row.pack(side="left")

        self.add_preset_var = ctk.StringVar(value=get_preset_labels()[0])
        self.add_preset_menu = ctk.CTkOptionMenu(
            add_row,
            values=get_preset_labels(),
            variable=self.add_preset_var,
            width=96,
            command=self._on_add_preset_change,
        )
        self.add_preset_menu.pack(side="left", padx=(0, 4))
        UITheme.style_option_menu(self.add_preset_menu)

        add_btn = ctk.CTkButton(
            add_row,
            text="+ 添加",
            width=72,
            command=self._add_annotation_mode,
        )
        add_btn.pack(side="left")
        UITheme.style_primary(add_btn)

        self._pending_add_preset = None

        self.annotation_list = ctk.CTkScrollableFrame(
            self._editor_card,
            height=150,
            fg_color="transparent",
            **UITheme.scrollable_frame_kwargs(),
        )
        self.annotation_list.pack(fill="x", padx=pad, pady=(0, pad))
        UITheme.style_scrollable_frame(self.annotation_list)

        SectionHeader(self._editor_card, text="批注内容").pack(
            fill="x", padx=pad, pady=(pad, 4)
        )

        style_panel = ctk.CTkFrame(self._editor_card, fg_color=UITheme.PURPLE_50, corner_radius=8)
        style_panel.pack(fill="x", padx=pad, pady=(0, 6))

        from src.utils.system_fonts import get_system_font_families

        self._system_fonts = get_system_font_families()
        default_font = self.settings.annotation.style.font_family or "Microsoft YaHei"
        if default_font not in self._system_fonts:
            self._system_fonts.insert(0, default_font)

        font_row = ctk.CTkFrame(style_panel, fg_color="transparent")
        font_row.pack(fill="x", padx=8, pady=(6, 4))
        ctk.CTkLabel(
            font_row, text="字体", width=36, anchor="w",
            font=UITheme.font_caption(), text_color=UITheme.TEXT_MUTED,
        ).pack(side="left")
        self.font_family_var = ctk.StringVar(value=default_font)
        self.font_family_menu = ctk.CTkComboBox(
            font_row,
            values=self._system_fonts,
            variable=self.font_family_var,
            width=200,
            command=self._on_font_family_change,
        )
        self.font_family_menu.pack(side="left", fill="x", expand=True, padx=(4, 0))
        UITheme.style_combo(self.font_family_menu)

        size_row = ctk.CTkFrame(style_panel, fg_color="transparent")
        size_row.pack(fill="x", padx=8, pady=4)
        ctk.CTkLabel(
            size_row, text="字号", width=36, anchor="w",
            font=UITheme.font_caption(), text_color=UITheme.TEXT_MUTED,
        ).pack(side="left")
        self.font_size_var = ctk.StringVar(value="12")
        self.font_size_entry = ctk.CTkEntry(
            size_row, width=52, textvariable=self.font_size_var,
            placeholder_text="6–96",
        )
        self.font_size_entry.pack(side="left", padx=(4, 6))
        UITheme.style_entry(self.font_size_entry)
        self.font_size_entry.bind("<Return>", lambda _e: self._on_font_size_apply())
        self.font_size_entry.bind("<FocusOut>", lambda _e: self._on_font_size_apply())

        ctk.CTkLabel(
            size_row, text="pt", font=UITheme.font_caption(), text_color=UITheme.TEXT_MUTED,
        ).pack(side="left")

        orient_row = ctk.CTkFrame(style_panel, fg_color="transparent")
        orient_row.pack(fill="x", padx=8, pady=4)
        ctk.CTkLabel(
            orient_row, text="方向", width=36, anchor="w",
            font=UITheme.font_caption(), text_color=UITheme.TEXT_MUTED,
        ).pack(side="left")
        from src.utils.text_orientation import ORIENTATION_LABELS

        self.text_orientation_var = ctk.StringVar(value=ORIENTATION_LABELS["horizontal"])
        self.text_orientation_menu = ctk.CTkSegmentedButton(
            orient_row,
            values=[ORIENTATION_LABELS["horizontal"], ORIENTATION_LABELS["vertical"]],
            variable=self.text_orientation_var,
            command=self._on_text_orientation_change,
        )
        self.text_orientation_menu.pack(side="left", fill="x", expand=True, padx=(4, 0))
        UITheme.style_segmented_panel(self.text_orientation_menu)

        color_row = ctk.CTkFrame(style_panel, fg_color="transparent")
        color_row.pack(fill="x", padx=8, pady=(4, 8))
        ctk.CTkLabel(
            color_row, text="颜色", width=36, anchor="w",
            font=UITheme.font_caption(), text_color=UITheme.TEXT_MUTED,
        ).pack(side="left", padx=(0, 4))

        self.color_var = ctk.StringVar(value=UITheme.ANNOTATION_COLOR_DEFAULT)
        for color in UITheme.ANNOTATION_COLORS:
            ctk.CTkButton(
                color_row,
                text="",
                width=24,
                height=24,
                corner_radius=12,
                command=lambda c=color: self._set_color(c),
                **UITheme.annotation_color_swatch_params(color),
            ).pack(side="left", padx=2)
        self.color_pick_btn = ctk.CTkButton(
            color_row,
            text="调色盘",
            width=56,
            height=24,
            command=self._pick_annotation_color,
        )
        self.color_pick_btn.pack(side="left", padx=(6, 0))
        UITheme.style_secondary(self.color_pick_btn)

        self.annotation_input = ctk.CTkTextbox(
            self._editor_card,
            height=200,
            wrap="word",
            activate_scrollbars=True,
            font=(default_font, 12),
        )
        self.annotation_input.pack(fill="both", expand=True, padx=pad, pady=4)
        UITheme.style_textbox(self.annotation_input)

        btn_frame = ctk.CTkFrame(self._editor_card, fg_color="transparent")
        btn_frame.pack(fill="x", padx=pad, pady=pad)

        self.save_btn = ctk.CTkButton(
            btn_frame,
            text="保存",
            width=80,
            command=self._save_annotation
        )
        self.save_btn.pack(side="left", padx=5)

        self.delete_btn = ctk.CTkButton(
            btn_frame,
            text="删除",
            width=80,
            command=self._delete_annotation
        )
        self.delete_btn.pack(side="left", padx=5)

        self._on_add_preset_change(self.add_preset_var.get())

        self.mode_hint = ctk.CTkLabel(
            self.sidebar_frame,
            text="滚轮滚动 · ± 缩放 · 双击译文编辑 · 可拖动微调位置",
        )
        self.mode_hint.pack(pady=(0, UITheme.PAD))

    def _apply_theme(self) -> None:
        """应用设计系统"""
        UITheme.style_card(self.file_list_frame)
        UITheme.style_card(self.preview_frame)
        UITheme.style_card(self.sidebar_frame)
        if hasattr(self, "_editor_card"):
            UITheme.style_card(self._editor_card, elevated=True)
        UITheme.style_primary(self.save_btn)
        UITheme.style_soft_danger(self.delete_btn)
        if hasattr(self, "add_preset_menu"):
            UITheme.style_option_menu(self.add_preset_menu)
        if hasattr(self, "font_family_menu"):
            UITheme.style_combo(self.font_family_menu)
        if hasattr(self, "text_orientation_menu"):
            UITheme.style_segmented_panel(self.text_orientation_menu)
        if hasattr(self, "color_pick_btn"):
            UITheme.style_secondary(self.color_pick_btn)
        if hasattr(self, "font_size_entry"):
            UITheme.style_entry(self.font_size_entry)
        UITheme.muted_label(self.file_hint)
        UITheme.muted_label(self.mode_hint)
        if hasattr(self, "_nav_inner"):
            UITheme.style_nav_bar(self._nav_inner)
            for btn in getattr(self, "_nav_buttons", []):
                UITheme.style_nav_button(btn)
            UITheme.title_label(self.page_label)

    def _configure_layout(self) -> None:
        """配置布局（状态栏先贴底，避免被 expand 的中间区域挤出屏幕）"""
        p = UITheme.PAD
        self.status_bar.pack(side="bottom", fill="x", padx=p, pady=(0, p))
        self.toolbar.pack(side="top", fill="x", padx=p, pady=(p, UITheme.PAD_SM))
        self.content_frame.pack(
            side="top", fill="both", expand=True, padx=p, pady=(UITheme.PAD_SM, UITheme.PAD_SM)
        )
        self.file_list_frame.pack(side="left", fill="y", padx=(0, UITheme.PAD_SM))
        self.preview_frame.pack(side="left", fill="both", expand=True, padx=UITheme.PAD_SM)
        self.sidebar_frame.pack(side="right", fill="y", padx=(UITheme.PAD_SM, 0))

    def _current_file_path(self) -> str:
        if 0 <= self.current_file_index < len(self.selected_files):
            return self.selected_files[self.current_file_index]
        return ""

    def get_render_pdf_path(self, file_path: str = "") -> str:
        """当前文件对应的 PDF 路径（PPT 为转换后的 PDF）"""
        from src.utils.ppt_converter import get_render_pdf_path

        path = file_path or self._current_file_path()
        return get_render_pdf_path(path, self.converted_pdf_paths)

    def _canvas_scale(self) -> float:
        """PDF 点坐标 → 画布像素的缩放比"""
        return self.zoom_level * 2

    def _inline_font_pixels(self, marker: AnnotationMarker) -> int:
        """原位译文在画布上的像素字号（与页面缩放后的 PDF 点一致）。"""
        from src.utils.block_font_size import GENERATED_INLINE_FONT_PT

        scale = self._canvas_scale()
        pt = float(getattr(marker, "font_size", 12) or 12)
        if (getattr(marker, "original_text", "") or "").strip():
            pt = float(GENERATED_INLINE_FONT_PT)
        return max(8, min(32, int(round(pt * scale))))

    def _marker_to_canvas(self, marker: AnnotationMarker) -> Tuple[float, float]:
        scale = self._canvas_scale()
        return (
            marker.x * scale + self._page_offset_x,
            marker.y * scale + self._page_offset_y,
        )

    def _canvas_to_pdf(self, cx: float, cy: float) -> Tuple[float, float]:
        scale = self._canvas_scale()
        if scale <= 0:
            return cx, cy
        return (
            (cx - self._page_offset_x) / scale,
            (cy - self._page_offset_y) / scale,
        )

    def _normalize_markers(self, page_num: int) -> None:
        """将旧版画布像素坐标转为 PDF 点坐标，并限制在页面内。

        批注应始终以 PDF 点存储；显示时由 _marker_to_canvas 按当前 zoom 换算。
        勿在每次缩放重绘时反复换算，否则会在缩小/放大时产生漂移。
        """
        if not self.pdf_doc:
            return
        markers = self.annotations.get(page_num, [])
        if not markers:
            return

        page = self.pdf_doc[page_num]
        pw, ph = page.rect.width, page.rect.height
        scale = self._canvas_scale()
        if scale <= 0:
            return
        size = AnnotationMarker.MARKER_SIZE

        for marker in markers:
            x, y = float(marker.x), float(marker.y)
            if x > pw:
                x /= scale
            if y > ph:
                y /= scale
            if getattr(marker, "display_mode", "marker") == "inline":
                marker.x = int(max(0, min(x, pw - 1)))
                marker.y = int(max(0, min(y, ph - 1)))
                from src.utils.block_font_size import (
                    GENERATED_INLINE_FONT_PT,
                    clamp_font_size_pt,
                )

                if (getattr(marker, "original_text", "") or "").strip():
                    marker.font_size = GENERATED_INLINE_FONT_PT
                elif getattr(marker, "font_size", 0):
                    marker.font_size = clamp_font_size_pt(marker.font_size)
                continue
            marker.x = int(max(0, min(x, pw - size)))
            marker.y = int(max(0, min(y, ph - size)))

    def _file_key(self, path: str) -> str:
        from src.utils.file_utils import file_key

        return file_key(path)

    def _rekey_annotations_by_file(self) -> None:
        """统一批注存储键，合并同一路径的重复项"""
        merged: Dict[str, Dict[int, List[AnnotationMarker]]] = {}
        for path, pages in self.annotations_by_file.items():
            key = self._file_key(path)
            bucket = merged.setdefault(key, {})
            for page, markers in pages.items():
                bucket[int(page)] = list(markers)
        self.annotations_by_file = merged

    def _get_stored_annotations(self, file_path: str) -> Dict[int, List[AnnotationMarker]]:
        key = self._file_key(file_path)
        if key in self.annotations_by_file:
            return self.annotations_by_file[key]
        base = os.path.basename(file_path)
        for path, pages in self.annotations_by_file.items():
            if os.path.basename(path) == base:
                return pages
        return {}

    def _persist_current_file_annotations(self) -> None:
        path = self._current_file_path()
        if path:
            key = self._file_key(path)
            self.annotations_by_file[key] = {
                page: list(markers) for page, markers in self.annotations.items()
            }

    def _restore_file_annotations(self, file_path: str) -> None:
        stored = self._get_stored_annotations(file_path)
        self.annotations = {
            int(page): list(markers) for page, markers in stored.items()
        }
        self.selected_marker = None
        if self.pdf_doc:
            for page_num in list(self.annotations.keys()):
                if 0 <= page_num < self.total_pages:
                    self._normalize_markers(page_num)

    def schedule_persist(self) -> None:
        if not self.settings.app.auto_save:
            return
        if self._autosave_job:
            self.after_cancel(self._autosave_job)
        self._autosave_job = self.after(800, self._do_persist)

    def _flush_persist(self) -> None:
        """立即保存会话（不等待防抖）"""
        if self._autosave_job:
            self.after_cancel(self._autosave_job)
            self._autosave_job = None
        from src.services.project_service import save_session

        save_session(self)

    def _do_persist(self) -> None:
        self._autosave_job = None
        self._flush_persist()

    def _restore_last_session(self) -> None:
        from src.services.project_service import restore_session

        if restore_session(self):
            self._rekey_annotations_by_file()

    def _on_close(self) -> None:
        self._persist_current_file_annotations()
        self._flush_persist()
        if self.pdf_doc:
            self.pdf_doc.close()
        self.destroy()

    def update_status(self, message: str) -> None:
        """更新状态栏消息"""
        self.status_bar.set_message(message)

    def update_progress(self, current: int, total: int, status: str) -> None:
        """更新进度"""
        self.status_bar.set_progress(current, total, status)

    def sync_web_preview(self) -> None:
        """同步批注数据到 PDF.js 预览服务"""
        if hasattr(self, "web_preview"):
            self.web_preview.update_from_app(self)
        self.schedule_persist()

    def _current_preview_ink_key(self) -> str:
        from src.utils.file_utils import file_key

        path = self._current_file_path()
        return file_key(path) if path else ""

    def _capture_undo_snapshot(self) -> dict:
        """当前文件批注 + 预览墨迹的快照。"""
        from src.services.project_service import _marker_to_dict

        fk = self._current_preview_ink_key()
        if not fk:
            return {}

        ann_pages = self.annotations_by_file.get(fk, {})
        annotations = {
            str(p): [_marker_to_dict(m) for m in markers]
            for p, markers in ann_pages.items()
        }
        ink_pages = {
            str(p): list(strokes)
            for p, strokes in (self.preview_ink_by_file.get(fk, {}) or {}).items()
        }
        sel = self.selected_marker
        selected = None
        if sel and self.current_page in self.annotations:
            try:
                selected = self.annotations[self.current_page].index(sel)
            except ValueError:
                selected = None

        return {
            "file_key": fk,
            "page": self.current_page,
            "annotations": annotations,
            "ink": ink_pages,
            "selected_index": selected,
        }

    def _record_undo_before(self) -> None:
        if self._undo_restoring:
            return
        snap = self._capture_undo_snapshot()
        if snap:
            self._undo_stack.push(snap)

    def _flush_drag_undo_snap(self) -> None:
        if self._undo_restoring:
            self._drag_undo_snap = None
            return
        if self._drag_undo_snap:
            self._undo_stack.push(self._drag_undo_snap)
        self._drag_undo_snap = None

    def _restore_undo_snapshot(self, snap: dict) -> None:
        from src.services.project_service import _dict_to_marker

        if not snap:
            return
        fk = snap.get("file_key")
        if not fk:
            return

        self._undo_restoring = True
        try:
            pages: Dict[int, List[AnnotationMarker]] = {}
            for p_str, items in (snap.get("annotations") or {}).items():
                pages[int(p_str)] = [_dict_to_marker(d) for d in items]
            self.annotations_by_file[fk] = pages

            ink: Dict[int, List[dict]] = {}
            for p_str, strokes in (snap.get("ink") or {}).items():
                ink[int(p_str)] = list(strokes)
            self.preview_ink_by_file[fk] = ink

            page = int(snap.get("page", 0))
            if 0 <= page < self.total_pages:
                self.current_page = page
                if self.pdf_doc:
                    self._render_page()

            self.annotations = {
                int(p): list(markers)
                for p, markers in self.annotations_by_file.get(fk, {}).items()
            }
            self._persist_current_file_annotations()

            idx = snap.get("selected_index")
            self.selected_marker = None
            markers = self.annotations.get(self.current_page, [])
            if idx is not None and 0 <= idx < len(markers):
                self._select_marker(markers[idx])
            else:
                self.annotation_input.delete("1.0", "end")
                self._update_annotation_list()

            self._show_annotations()
            self.sync_web_preview()
            self.update_status("已撤销上一步操作")
        finally:
            self._undo_restoring = False

    def _on_undo(self, event=None) -> Optional[str]:
        if self._is_text_entry_focused():
            return None
        if self._inline_edit_marker:
            self._cancel_inline_text_edit()
        fk = self._current_preview_ink_key()
        snap = None
        while self._undo_stack.can_undo:
            candidate = self._undo_stack.pop()
            if not candidate:
                continue
            if not fk or candidate.get("file_key") == fk:
                snap = candidate
                break
        if not snap:
            self.update_status("没有可撤销的操作")
            return "break"
        self._restore_undo_snapshot(snap)
        return "break"

    def save_preview_ink_to_document(self, *, overwrite: bool = False) -> str:
        """将预览中的笔 / 荧光笔墨迹写入 PDF（PPT 写入转换后的 PDF 副本）。"""
        import fitz
        from datetime import datetime

        from src.utils.file_utils import file_key
        from src.utils.pdf_annotation import draw_page_annotations
        from src.utils.pdf_ink import draw_page_ink
        from src.utils.pdf_save import save_fitz_document
        from src.utils.preview_ink_store import normalize_ink_pages

        source = self._current_file_path()
        if not source or not self.pdf_doc:
            raise ValueError("请先在桌面应用中打开 PDF 或 PPT")

        key = file_key(source)
        ink_pages = normalize_ink_pages(self.preview_ink_by_file.get(key, {}))
        if not any(ink_pages.values()):
            raise ValueError("没有可保存的笔记，请先用笔或荧光笔在预览中绘制")

        render_pdf = self.get_render_pdf_path(source)
        doc = fitz.open(render_pdf)
        markers_by_page = self._get_stored_annotations(source) or self.annotations

        for page_num in range(doc.page_count):
            strokes = ink_pages.get(page_num, [])
            if strokes:
                draw_page_ink(doc[page_num], strokes)
            markers = markers_by_page.get(page_num, [])
            if markers:
                draw_page_annotations(doc[page_num], markers)

        ext = os.path.splitext(source)[1].lower()
        if overwrite and ext == ".pdf":
            out_path = source
        else:
            directory = os.path.dirname(source) or os.getcwd()
            stem = os.path.splitext(os.path.basename(source))[0]
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            label = "_笔记" if ext in (".ppt", ".pptx") else "_带笔记"
            out_path = os.path.join(directory, f"{stem}{label}_{ts}.pdf")

        save_fitz_document(doc, out_path, opened_from=render_pdf)

        if ext in (".ppt", ".pptx"):
            self.converted_pdf_paths[key] = out_path
            self._load_pdf(out_path, page=self.current_page, load_text_positions=False)
        elif overwrite and ext == ".pdf":
            self._load_pdf(source, page=self.current_page)

        if hasattr(self, "web_preview"):
            self.web_preview.state.pdf_path = out_path
            self.sync_web_preview()

        self.schedule_persist()
        return out_path

    def open_web_preview(self) -> None:
        """在浏览器中打开 PDF.js 批注预览"""
        import webbrowser

        self.sync_web_preview()
        webbrowser.open(self.web_preview.url)
        self._web_preview_opened = True
        self.update_status("预览已打开，点击黄色标记查看批注，点击空白处关闭")

    def update_file_list(self, files: List[str]) -> None:
        """更新文件列表显示"""
        # 清空现有列表
        for widget in self.file_listbox.winfo_children():
            widget.destroy()

        if not files:
            self.file_hint.configure(text="导入文件")
            return

        self.file_hint.configure(text="")

        for idx, file_path in enumerate(files):
            file_name = os.path.basename(file_path)
            lower = file_name.lower()
            if lower.endswith(".pdf"):
                icon = "📄"
            elif lower.endswith((".ppt", ".pptx")):
                icon = "📊"
            else:
                icon = "📁"

            is_current = idx == self.current_file_index
            display_name = UITheme.truncate_filename(file_name)

            row = ctk.CTkFrame(
                self.file_listbox,
                fg_color=UITheme.PURPLE_100 if is_current else UITheme.SURFACE,
                border_color=UITheme.PURPLE_400 if is_current else UITheme.BORDER,
                border_width=1 if is_current else 0,
                corner_radius=UITheme.RADIUS,
            )
            row.pack(fill="x", pady=3, padx=2)
            row.grid_columnconfigure(0, weight=1)

            name_label = ctk.CTkLabel(
                row,
                text=f"{icon} {display_name}",
                anchor="w",
                justify="left",
                font=UITheme.font_body(),
                text_color=UITheme.PURPLE_800 if is_current else UITheme.TEXT,
                cursor="hand2",
            )
            name_label.grid(row=0, column=0, sticky="ew", padx=(10, 4), pady=8)

            def select_file(_event=None, i=idx):
                self._on_file_select(i)

            name_label.bind("<Button-1>", select_file)
            row.bind("<Button-1>", select_file)
            UITheme.bind_tooltip(name_label, file_name)
            UITheme.bind_tooltip(row, file_name)

            remove_btn = ctk.CTkButton(
                row,
                text="删",
                width=32,
                height=28,
                command=lambda i=idx: self._remove_file(i),
            )
            remove_btn.grid(row=0, column=1, padx=(0, 6), pady=6)
            UITheme.style_icon_dismiss(remove_btn)

        # 默认选中第一个文件
        if files and self.current_file_index < 0:
            self._on_file_select(0)

    def _on_file_select(self, index: int, page: int = None) -> None:
        """选择文件（各文件批注独立存储）"""
        if 0 <= index < len(self.selected_files):
            self._persist_current_file_annotations()
            self._undo_stack.clear()
            self._drag_undo_snap = None
            self.current_file_index = index
            file_path = self.selected_files[index]
            file_name = os.path.basename(file_path)

            self._load_file_content(file_path, page=page)
            self._restore_file_annotations(file_path)
            self._show_annotations()
            self._update_annotation_list()

            self.update_status(f"已选择: {file_name}")
            self.update_file_list(self.selected_files)
            self.sync_web_preview()

    def _load_file_content(self, file_path: str, page: int = None) -> None:
        """加载文件内容"""
        try:
            if file_path.lower().endswith('.pdf'):
                self._load_pdf(file_path, page=page)
            elif file_path.lower().endswith(('.ppt', '.pptx')):
                self._load_ppt(file_path, page=page)
            else:
                show_warning(self, "警告", "不支持的文件格式")
        except Exception as e:
            show_warning(self, "错误", f"加载文件时出错: {str(e)}")

    def _get_page_text(self, page_num: int) -> str:
        """获取当前文件指定页的文字"""
        from src.utils.page_image import extract_page_text_for_annotation

        return extract_page_text_for_annotation(
            page_num,
            pdf_doc=self.pdf_doc,
            source_path=self._current_file_path(),
        )

    def _load_pdf(self, file_path: str, page: int = None, load_text_positions: bool = True) -> None:
        """加载PDF文件并渲染为图片"""
        self.ocr_text_positions = {}
        self.ocr_text_sources = {}
        try:
            import fitz  # PyMuPDF

            if self.pdf_doc:
                self.pdf_doc.close()

            self.pdf_doc = fitz.open(file_path)
            self.total_pages = len(self.pdf_doc)
            if page is None:
                self.current_page = 0
            else:
                self.current_page = max(0, min(page, self.total_pages - 1))

            if load_text_positions:
                self._load_text_positions(file_path)
            self._render_page()

            self.update_status(f"已加载PDF: {self.total_pages}页")

        except ImportError:
            show_warning(self, "错误", "请安装PyMuPDF: pip install pymupdf")
        except Exception as e:
            show_warning(self, "错误", f"读取PDF出错: {str(e)}")

    def _load_text_positions(self, file_path: str) -> None:
        """使用 LiteParse 加载文本位置信息"""
        try:
            from liteparse import LiteParse
            parser = LiteParse(ocr_enabled=False)  # 不需要 OCR，速度更快
            result = parser.parse(file_path)

            # 存储每页的文本位置信息
            self.text_positions = {}
            for page_data in result.pages:
                page_num = page_data.page_num - 1  # LiteParse 是 1-indexed
                positions = []

                if hasattr(page_data, 'text_items') and page_data.text_items:
                    for item in page_data.text_items:
                        if hasattr(item, 'text') and item.text.strip():
                            h = float(getattr(item, 'height', 0) or 0)
                            pos = {
                                "text": item.text,
                                "x": getattr(item, 'x', 0),
                                "y": getattr(item, 'y', 0),
                                "width": getattr(item, 'width', 0),
                                "height": h,
                            }
                            fs = getattr(item, "font_size", None) or getattr(item, "size", None)
                            if fs is not None:
                                pos["font_size"] = fs
                            elif h >= 6:
                                from src.utils.block_font_size import font_size_from_line_height
                                pos["font_size"] = font_size_from_line_height(h)
                            positions.append(pos)

                self.text_positions[page_num] = positions

            print(f"LiteParse: 已加载 {len(self.text_positions)} 页的文本位置信息")

        except ImportError:
            print("提示: 未安装 liteparse，将使用基本模式")
            self.text_positions = {}
        except Exception as e:
            print(f"LiteParse 加载失败: {e}")
            self.text_positions = {}

    def _load_pptx_text_positions(self, file_path: str) -> None:
        """从 PPTX 提取各页文字位置信息"""
        try:
            from pptx import Presentation

            from src.utils.block_font_size import font_size_from_pptx_shape

            prs = Presentation(file_path)
            self.ppt_slide_emu = (int(prs.slide_width), int(prs.slide_height))
            self.text_positions = {}
            for slide_index, slide in enumerate(prs.slides):
                positions = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            pos = {
                                "text": text,
                                "x": getattr(shape, "left", 0),
                                "y": getattr(shape, "top", 0),
                                "width": getattr(shape, "width", 0),
                                "height": getattr(shape, "height", 0),
                            }
                            fs = font_size_from_pptx_shape(shape)
                            if fs is not None:
                                pos["font_size"] = fs
                            positions.append(pos)
                self.text_positions[slide_index] = positions
        except Exception as e:
            print(f"PPTX 文字提取失败: {e}")
            self.text_positions = {}

    def _load_ppt(self, file_path: str, page: int = None) -> None:
        """加载 PPT/PPTX（转换为 PDF 后预览，批注时用幻灯片图片）"""
        try:
            from src.utils.ppt_converter import convert_ppt_to_pdf

            if file_path.lower().endswith(".pptx"):
                self._load_pptx_text_positions(file_path)

            pdf_path = convert_ppt_to_pdf(file_path)
            self.converted_pdf_paths[self._file_key(file_path)] = pdf_path
            self.zoom_level = self.PPT_IMPORT_ZOOM
            self._load_pdf(pdf_path, page=page, load_text_positions=False)
            self.update_status(f"已加载 PPT: {self.total_pages} 页")
        except Exception as e:
            show_warning(self, "错误", f"无法打开 PPT：{e}")
            raise

    def _ensure_preview_shell(self) -> None:
        """固定预览区：底部缩放/翻页栏 + 带滚动条的画布（PDF/PPT 共用）"""
        if self._preview_shell_ready:
            return

        nav_outer = ctk.CTkFrame(self.preview_frame, fg_color="transparent")
        nav_outer.pack(side="bottom", fill="x", padx=UITheme.PAD, pady=UITheme.PAD)

        self._build_ink_toolbar(nav_outer)

        self._nav_inner = ctk.CTkFrame(nav_outer)
        self._nav_inner.pack(fill="x")
        UITheme.style_nav_bar(self._nav_inner)

        nav = self._nav_inner
        self._nav_buttons = []

        zoom_frame = ctk.CTkFrame(nav, fg_color="transparent")
        zoom_frame.pack(side="left", padx=UITheme.PAD, pady=UITheme.PAD_SM)

        zoom_out_btn = ctk.CTkButton(zoom_frame, text="−", width=36, command=self._zoom_out)
        zoom_out_btn.pack(side="left", padx=2)
        UITheme.style_nav_button(zoom_out_btn)
        self._nav_buttons.append(zoom_out_btn)

        self.zoom_label = ctk.CTkLabel(
            zoom_frame,
            text=f"{int(self.zoom_level * 100)}%",
            width=56,
            font=UITheme.font_section(),
            text_color=UITheme.PURPLE_800,
        )
        self.zoom_label.pack(side="left", padx=6)

        zoom_in_btn = ctk.CTkButton(zoom_frame, text="+", width=36, command=self._zoom_in)
        zoom_in_btn.pack(side="left", padx=2)
        UITheme.style_nav_button(zoom_in_btn)
        self._nav_buttons.append(zoom_in_btn)

        self.prev_btn = ctk.CTkButton(nav, text="上一页", command=self._prev_page, width=96)
        self.prev_btn.pack(side="left", padx=UITheme.PAD)
        UITheme.style_nav_button(self.prev_btn)
        self._nav_buttons.append(self.prev_btn)

        self.page_label = ctk.CTkLabel(
            nav,
            text="0 / 0",
            font=UITheme.font_section(),
            text_color=UITheme.TEXT,
        )
        self.page_label.pack(side="left", expand=True)

        self.next_btn = ctk.CTkButton(nav, text="下一页", command=self._next_page, width=96)
        self.next_btn.pack(side="right", padx=UITheme.PAD)
        UITheme.style_nav_button(self.next_btn)
        self._nav_buttons.append(self.next_btn)

        self._viewer_host = ctk.CTkFrame(self.preview_frame, fg_color="transparent")
        self._viewer_host.pack(side="top", fill="both", expand=True, padx=5, pady=(5, 0))

        self._canvas_panel = tk.Frame(self._viewer_host, bg=UITheme.BG, highlightthickness=0)
        self._canvas_panel.pack(fill="both", expand=True)
        self._canvas_panel.grid_rowconfigure(0, weight=1)
        self._canvas_panel.grid_columnconfigure(0, weight=1)

        self.canvas = tk.Canvas(
            self._canvas_panel,
            bg=UITheme.SURFACE,
            highlightthickness=0,
        )
        self._v_scroll = ctk.CTkScrollbar(
            self._canvas_panel,
            orientation="vertical",
            command=self.canvas.yview,
        )
        self._h_scroll = ctk.CTkScrollbar(
            self._canvas_panel,
            orientation="horizontal",
            command=self.canvas.xview,
        )
        self.canvas.configure(
            yscrollcommand=self._v_scroll.set,
            xscrollcommand=self._h_scroll.set,
        )
        UITheme.style_scrollbar(self._v_scroll)
        UITheme.style_scrollbar(self._h_scroll)

        self.canvas.grid(row=0, column=0, sticky="nsew")
        self._v_scroll.grid(row=0, column=1, sticky="ns")
        self._h_scroll.grid(row=1, column=0, sticky="ew")
        for sb in (self._v_scroll, self._h_scroll):
            sb.bind("<ButtonRelease-1>", lambda _e: self._schedule_viewport_guide())

        self._bind_canvas_events()
        self._preview_shell_ready = True

    def _build_ink_toolbar(self, parent: ctk.CTkFrame) -> None:
        """与 Web 预览一致的笔 / 荧光笔 / 橡皮擦。"""
        from src.ui.canvas_ink import INK_COLORS

        bar = ctk.CTkFrame(parent, fg_color=UITheme.SURFACE, corner_radius=UITheme.RADIUS)
        bar.pack(fill="x", pady=(0, 6))
        UITheme.style_card(bar, elevated=False)

        inner = ctk.CTkFrame(bar, fg_color="transparent")
        inner.pack(fill="x", padx=8, pady=6)

        ctk.CTkLabel(
            inner,
            text="绘图",
            font=UITheme.font_caption(),
            text_color=UITheme.TEXT_MUTED,
        ).pack(side="left", padx=(0, 8))

        tools = [
            ("select", "选择"),
            ("pen", "笔"),
            ("highlighter", "荧光笔"),
            ("eraser", "橡皮"),
        ]
        self._ink_tool_buttons.clear()
        for tool_id, label in tools:
            btn = ctk.CTkButton(
                inner,
                text=label,
                width=56,
                height=28,
                command=lambda t=tool_id: self._set_canvas_tool(t),
            )
            btn.pack(side="left", padx=2)
            self._ink_tool_buttons[tool_id] = btn

        clear_btn = ctk.CTkButton(
            inner,
            text="清除本页",
            width=72,
            height=28,
            command=self._clear_page_ink,
        )
        clear_btn.pack(side="left", padx=(8, 4))
        UITheme.style_soft_danger(clear_btn)

        color_row = ctk.CTkFrame(inner, fg_color="transparent")
        color_row.pack(side="right", padx=(8, 0))
        for color in INK_COLORS[:8]:
            ctk.CTkButton(
                color_row,
                text="",
                width=20,
                height=20,
                corner_radius=10,
                fg_color=color,
                hover_color=color,
                command=lambda c=color: self._set_ink_color(c),
            ).pack(side="left", padx=2)

        self._set_canvas_tool("select")

    def _set_canvas_tool(self, tool: str) -> None:
        self.canvas_tool = tool
        self._current_ink_stroke = None
        for tid, btn in self._ink_tool_buttons.items():
            if tid == tool:
                UITheme.style_primary(btn)
            else:
                UITheme.style_secondary(btn)
        hints = {
            "select": "选择：拖动批注；滚轮滚动页面",
            "pen": "笔：在页面上手绘，墨迹会同步到 Web 预览",
            "highlighter": "荧光笔：宽笔触高亮",
            "eraser": "橡皮：擦除笔迹与荧光笔",
        }
        if hasattr(self, "mode_hint"):
            self.mode_hint.configure(text=hints.get(tool, self.mode_hint.cget("text")))

    def _set_ink_color(self, color: str) -> None:
        self.ink_color = color

    def _page_ink_bucket(self) -> List[dict]:
        key = self._current_preview_ink_key()
        if not key:
            return []
        bucket = self.preview_ink_by_file.setdefault(key, {})
        return bucket.setdefault(self.current_page, [])

    def _page_ink_strokes(self) -> List[dict]:
        from src.utils.preview_ink_store import normalize_ink_pages

        key = self._current_preview_ink_key()
        if not key:
            return []
        pages = normalize_ink_pages(self.preview_ink_by_file.get(key, {}))
        return list(pages.get(self.current_page, []))

    def _persist_page_ink(self) -> None:
        from src.utils.preview_ink_store import normalize_ink_pages

        key = self._current_preview_ink_key()
        if not key:
            return
        strokes = self._page_ink_bucket()
        pages = normalize_ink_pages(self.preview_ink_by_file.get(key, {}))
        if strokes:
            pages[self.current_page] = strokes
        elif self.current_page in pages:
            del pages[self.current_page]
        self.preview_ink_by_file[key] = pages
        self.schedule_persist()
        self.sync_web_preview()

    def _redraw_page_ink(self) -> None:
        if not self.canvas or not self.page_image:
            return
        from src.ui.canvas_ink import redraw_page_ink

        pw, ph = float(self.page_image.width), float(self.page_image.height)
        redraw_page_ink(
            self.canvas,
            self._page_ink_strokes(),
            float(self._page_offset_x),
            float(self._page_offset_y),
            pw,
            ph,
        )

    def _canvas_norm_point(self, cx: float, cy: float) -> Dict[str, float]:
        from src.ui.canvas_ink import canvas_to_norm

        if not self.page_image:
            return {"x": 0.0, "y": 0.0}
        return canvas_to_norm(
            cx,
            cy,
            float(self._page_offset_x),
            float(self._page_offset_y),
            float(self.page_image.width),
            float(self.page_image.height),
        )

    def _start_ink_stroke(self, cx: float, cy: float) -> None:
        from src.ui.canvas_ink import (
            append_stroke_point,
            new_stroke,
            stroke_width_for_tool,
        )

        if not self.page_image or self.canvas_tool not in ("pen", "highlighter"):
            return
        norm = self._canvas_norm_point(cx, cy)
        pw = float(self.page_image.width)
        width = stroke_width_for_tool(self.canvas_tool, pw)
        self._current_ink_stroke = new_stroke(
            self.canvas_tool,
            self.ink_color,
            width,
            [norm],
            page_width_px=pw,
        )

    def _extend_ink_stroke(self, cx: float, cy: float) -> None:
        from src.ui.canvas_ink import append_stroke_point, draw_stroke_on_canvas

        stroke = self._current_ink_stroke
        if not stroke or not self.page_image:
            return
        norm = self._canvas_norm_point(cx, cy)
        if not append_stroke_point(stroke, norm["x"], norm["y"]):
            return
        pw, ph = float(self.page_image.width), float(self.page_image.height)
        self.canvas.delete("ink_live")
        draw_stroke_on_canvas(
            self.canvas,
            stroke,
            float(self._page_offset_x),
            float(self._page_offset_y),
            pw,
            ph,
            tag="ink_live",
        )

    def _finish_ink_stroke(self) -> None:
        stroke = self._current_ink_stroke
        self._current_ink_stroke = None
        if not stroke or len(stroke.get("points") or []) < 2:
            if self.canvas:
                self.canvas.delete("ink_live")
            return
        self._page_ink_bucket().append(stroke)
        self._persist_page_ink()
        if self.canvas:
            self.canvas.delete("ink_live")
        self._redraw_page_ink()
        self.update_status("笔迹已保存（已同步预览）")

    def _erase_ink_at(self, cx: float, cy: float) -> None:
        from src.ui.canvas_ink import erase_strokes_at

        norm = self._canvas_norm_point(cx, cy)
        strokes = self._page_ink_strokes()
        new_strokes = erase_strokes_at(strokes, norm["x"], norm["y"])
        if len(new_strokes) == len(strokes):
            return
        key = self._current_preview_ink_key()
        if key:
            pages = self.preview_ink_by_file.setdefault(key, {})
            pages[self.current_page] = new_strokes
        self._persist_page_ink()
        self._redraw_page_ink()

    def _clear_page_ink(self) -> None:
        from src.ui.message_dialog import ask_yes_no

        if not self._page_ink_strokes():
            self.update_status("本页没有墨迹")
            return
        if not ask_yes_no(self, "清除墨迹", "确定清除当前页所有笔迹与荧光笔？"):
            return
        self._record_undo_before()
        key = self._current_preview_ink_key()
        if key:
            pages = self.preview_ink_by_file.setdefault(key, {})
            pages[self.current_page] = []
        self._persist_page_ink()
        self._redraw_page_ink()
        self.update_status("已清除本页墨迹")

    def _on_save_page_ink(self) -> None:
        try:
            path = self.save_preview_ink_to_document(overwrite=False)
            self.update_status(f"笔记已保存：{path}")
        except Exception as e:
            show_warning(self, "保存失败", str(e))

    def _bind_canvas_events(self) -> None:
        """画布交互：滚轮滚动（非缩放），批注拖拽等"""
        self.canvas.bind("<Button-1>", self._on_canvas_press)
        self.canvas.bind("<B1-Motion>", self._on_canvas_drag)
        self.canvas.bind("<ButtonRelease-1>", self._on_canvas_release)
        self.canvas.bind("<Double-Button-1>", self._on_canvas_double_click)

        for widget in (self.canvas, self._canvas_panel, self._viewer_host):
            widget.bind("<MouseWheel>", self._on_canvas_scroll)
            widget.bind("<Button-4>", self._on_canvas_scroll_linux_up)
            widget.bind("<Button-5>", self._on_canvas_scroll_linux_down)

        self.canvas.bind("<Configure>", self._on_canvas_configure)

    def _on_canvas_configure(self, event) -> None:
        """窗口尺寸变化时保持页面在可视区域居中"""
        if event.width < 2 or not self.page_image:
            return
        if self._canvas_configure_after:
            self.after_cancel(self._canvas_configure_after)
        self._canvas_configure_after = self.after(80, self._relayout_page_view)

    def _relayout_page_view(self) -> None:
        self._canvas_configure_after = None
        if not self.page_image or not self.canvas:
            return
        self._layout_page_in_canvas()
        self._show_annotations()

    def _layout_page_in_canvas(self) -> None:
        """将 PDF 页面置于滚动区域中央，并调整初始滚动位置"""
        if not self.page_image or not self.canvas:
            return

        pad = 24
        pw, ph = self.page_image.size
        self.canvas.update_idletasks()
        cw = max(self.canvas.winfo_width(), 1)
        ch = max(self.canvas.winfo_height(), 1)

        region_w = max(pw + pad * 2, cw)
        region_h = max(ph + pad * 2, ch)
        self._page_offset_x = (region_w - pw) // 2
        self._page_offset_y = (region_h - ph) // 2

        if self._page_image_id:
            self.canvas.coords(
                self._page_image_id,
                self._page_offset_x,
                self._page_offset_y,
            )
        else:
            self._page_image_id = self.canvas.create_image(
                self._page_offset_x,
                self._page_offset_y,
                anchor="nw",
                image=self.tk_image,
            )

        self.canvas.configure(scrollregion=(0, 0, region_w, region_h))

        if region_w > cw:
            self.canvas.xview_moveto(
                max(0, (self._page_offset_x + pw / 2 - cw / 2) / (region_w - cw))
            )
        else:
            self.canvas.xview_moveto(0)

        if region_h > ch:
            self.canvas.yview_moveto(
                max(0, (self._page_offset_y + ph / 2 - ch / 2) / (region_h - ch))
            )
        else:
            self.canvas.yview_moveto(0)

    def _on_canvas_scroll(self, event) -> str:
        """滚轮上下滚动页面（与 Web 预览一致，不用滚轮缩放）"""
        if not hasattr(self, "canvas"):
            return "break"
        delta = int(-1 * (event.delta / 120)) if event.delta else 0
        if delta:
            self.canvas.yview_scroll(delta, "units")
            self._schedule_viewport_guide()
        return "break"

    def _on_canvas_scroll_linux_up(self, event) -> str:
        self.canvas.yview_scroll(-3, "units")
        self._schedule_viewport_guide()
        return "break"

    def _on_canvas_scroll_linux_down(self, event) -> str:
        self.canvas.yview_scroll(3, "units")
        self._schedule_viewport_guide()
        return "break"

    def _destroy_preview_shell(self) -> None:
        for widget in self.preview_frame.winfo_children():
            widget.destroy()
        self._preview_shell_ready = False

    def _render_page(self) -> None:
        """渲染当前页面"""
        if not self.pdf_doc:
            return

        self._commit_inline_text_edit(silent=True)

        try:
            import fitz  # PyMuPDF
        except ImportError:
            show_warning(self, "错误", "请安装PyMuPDF: pip install pymupdf")
            return

        self._ensure_preview_shell()

        page = self.pdf_doc[self.current_page]
        mat = fitz.Matrix(self.zoom_level * 2, self.zoom_level * 2)
        pix = page.get_pixmap(matrix=mat)

        img_data = pix.tobytes("png")
        self.page_image = Image.open(io.BytesIO(img_data))
        self.tk_image = ImageTk.PhotoImage(self.page_image)

        self.canvas.delete("all")
        self._page_image_id = None
        self._page_image_id = self.canvas.create_image(0, 0, anchor="nw", image=self.tk_image)

        self._layout_page_in_canvas()
        self._show_annotations()

        bbox = self.canvas.bbox("all")
        if bbox:
            x1, y1, x2, y2 = bbox
            region = self.canvas.cget("scrollregion").split()
            if len(region) == 4:
                rw, rh = float(region[2]), float(region[3])
                self.canvas.configure(
                    scrollregion=(
                        min(0, x1),
                        min(0, y1),
                        max(rw, x2),
                        max(rh, y2),
                    )
                )

        self.page_label.configure(text=f"{self.current_page + 1} / {self.total_pages}")
        self.zoom_label.configure(text=f"{int(self.zoom_level * 100)}%")
        self._update_annotation_list()
        self.sync_web_preview()

    def _show_annotations(self) -> None:
        """显示当前页面的批注（小标记 + 双击弹层）"""
        if not self.canvas:
            return

        self.canvas.delete("annotation")
        self.canvas.delete("ink_live")
        self._redraw_page_ink()
        markers = self.annotations.get(self.current_page, [])
        canvas_w = (
            self._page_offset_x + self.page_image.width
            if self.page_image
            else 2000
        )

        for idx, marker in enumerate(markers, start=1):
            if getattr(marker, "display_mode", "marker") == "inline":
                if marker is self._inline_edit_marker:
                    continue
                self._draw_inline_annotation(marker)
                continue

            x, y = self._marker_to_canvas(marker)
            size = AnnotationMarker.MARKER_SIZE

            marker.icon_rect = (x, y, x + size, y + size)
            rect_id = self.canvas.create_rectangle(
                x, y, x + size, y + size,
                fill=UITheme.MARKER_FILL,
                outline=marker.color or UITheme.MARKER_OUTLINE_DEFAULT,
                width=2,
                tags="annotation",
            )
            self.canvas.create_text(
                x + size / 2, y + size / 2,
                text=str(idx),
                anchor="center",
                fill=UITheme.POPUP_TEXT,
                font=("Arial", 10, "bold"),
                tags="annotation",
            )
            marker.canvas_id = rect_id
            marker.text_id = rect_id

            if not marker.is_expanded:
                marker.popup_rect = None
                marker.close_rect = None
                continue

            pw = AnnotationMarker.POPUP_WIDTH
            ph = marker._calc_popup_height()
            px = x + AnnotationMarker.POPUP_OFFSET_X
            if px + pw > canvas_w:
                px = max(4, x - pw - 8)

            marker.popup_x = px
            marker.popup_y = y
            marker.popup_rect = (px, y, px + pw, y + ph)
            marker.close_rect = None

            self.canvas.create_rectangle(
                px + 3, y + 3, px + pw + 3, y + ph + 3,
                fill=UITheme.POPUP_SHADOW,
                outline="",
                tags="annotation",
            )

            popup_host = self._build_annotation_popup(marker, idx, pw, ph)
            marker.popup_window_id = self.canvas.create_window(
                px,
                y,
                window=popup_host,
                anchor="nw",
                width=pw,
                height=ph,
                tags="annotation",
            )

        try:
            self.canvas.tag_raise("annotation")
            self.canvas.tag_raise("viewport_guide")
        except tk.TclError:
            pass
        self._schedule_viewport_guide()

    def _schedule_viewport_guide(self) -> None:
        job = getattr(self, "_viewport_guide_after", None)
        if job:
            try:
                self.after_cancel(job)
            except tk.TclError:
                pass
        self._viewport_guide_after = self.after_idle(self._draw_viewport_visible_guide)

    def _draw_viewport_visible_guide(self) -> None:
        """主界面：虚线标出画布当前可见范围（非批注框）。"""
        self._viewport_guide_after = None
        if not self.canvas:
            return
        self.canvas.delete("viewport_guide")
        try:
            vw = max(int(self.canvas.winfo_width()), 1)
            vh = max(int(self.canvas.winfo_height()), 1)
        except tk.TclError:
            return
        sr = self.canvas.cget("scrollregion") or ""
        parts = sr.split()
        if len(parts) == 4:
            tw = float(parts[2]) - float(parts[0])
            th = float(parts[3]) - float(parts[1])
            if tw <= vw + 2 and th <= vh + 2:
                return
        x0 = self.canvas.canvasx(0)
        y0 = self.canvas.canvasy(0)
        x1 = self.canvas.canvasx(vw)
        y1 = self.canvas.canvasy(vh)
        self.canvas.create_rectangle(
            x0,
            y0,
            x1,
            y1,
            outline=UITheme.PURPLE_400,
            width=1,
            dash=(6, 4),
            tags="viewport_guide",
        )
        self.canvas.tag_raise("viewport_guide")

    def _draw_inline_annotation(self, marker: AnnotationMarker) -> None:
        """在原文旁直接绘制中文译文（可点击拖动微调位置）。"""
        cx, cy = self._marker_to_canvas(marker)
        scale = self._canvas_scale()
        font_px = self._inline_font_pixels(marker)
        family = self._marker_font_family(marker)
        placement = getattr(marker, "placement", "right") or "right"
        if placement == "above":
            anchor = "sw"
        elif placement == "right":
            anchor = "w"
        else:
            anchor = "nw"

        wrap = None
        if placement == "below" and marker.box_width:
            wrap = max(80, int(marker.box_width * scale))

        outline = marker.color or UITheme.PURPLE_800
        if marker == self.selected_marker:
            outline = UITheme.PURPLE_600

        marker.text_id = self.canvas.create_text(
            cx,
            cy,
            text=self._display_text_for_marker(marker),
            anchor=anchor,
            fill=outline,
            font=(family, -font_px),
            tags=("annotation", "annotation_inline"),
            width=wrap,
        )
        marker.canvas_id = marker.text_id
        self.canvas.update_idletasks()
        bbox = self.canvas.bbox(marker.text_id)
        pad = 6
        if bbox:
            marker.icon_rect = (
                bbox[0] - pad,
                bbox[1] - pad,
                bbox[2] + pad,
                bbox[3] + pad,
            )
        else:
            marker.icon_rect = (cx - pad, cy - pad, cx + pad, cy + pad)

    def _build_annotation_popup(
        self, marker: AnnotationMarker, idx: int, pw: int, ph: int
    ) -> ctk.CTkFrame:
        """固定尺寸弹层容器，正文在框内滚动，避免超出白框"""
        outline = marker.color or UITheme.MARKER_OUTLINE_DEFAULT

        def on_close() -> None:
            marker.collapse()
            self.selected_marker = None
            self._show_annotations()
            self._update_annotation_list()

        host = ctk.CTkFrame(
            self.canvas,
            width=pw,
            height=ph,
            fg_color=UITheme.POPUP_BG,
            border_color=outline,
            border_width=2,
            corner_radius=UITheme.RADIUS,
        )
        host.pack_propagate(False)
        host.grid_propagate(False)

        header = ctk.CTkFrame(host, fg_color="transparent", height=AnnotationMarker.TITLE_AREA)
        header.pack(fill="x", padx=8, pady=(6, 0))
        header.pack_propagate(False)

        ctk.CTkLabel(
            header,
            text=f"批注 {idx}",
            font=UITheme.font_section(),
            text_color=outline,
        ).pack(side="left", padx=4)

        close_btn = ctk.CTkButton(
            header,
            text="✕",
            width=28,
            height=28,
            corner_radius=UITheme.RADIUS_SM,
            fg_color=UITheme.PURPLE_100,
            hover_color=UITheme.PURPLE_200,
            text_color=outline,
            command=on_close,
        )
        close_btn.pack(side="right", padx=2)
        marker.close_rect = None

        body_h = max(
            ph - AnnotationMarker.TITLE_AREA - AnnotationMarker.POPUP_PADDING,
            48,
        )
        text_box = ctk.CTkTextbox(
            host,
            height=body_h,
            wrap="word",
            font=UITheme.font_body(),
            fg_color=UITheme.POPUP_BG,
            text_color=UITheme.POPUP_TEXT,
            border_width=0,
            activate_scrollbars=True,
        )
        text_box.pack(fill="both", expand=True, padx=10, pady=(2, 10))
        UITheme.style_textbox(text_box)
        text_box.configure(fg_color=UITheme.POPUP_BG, border_width=0)
        text_box.insert("1.0", marker.text)
        text_box.configure(state="disabled")

        return host

    def _point_in_rect(self, x: float, y: float, rect) -> bool:
        if not rect:
            return False
        x1, y1, x2, y2 = rect
        return x1 <= x <= x2 and y1 <= y <= y2

    def _find_marker_at(self, x: float, y: float, *, icon_only: bool = False) -> AnnotationMarker:
        """根据坐标查找批注"""
        markers = self.annotations.get(self.current_page, [])
        for marker in reversed(markers):
            if self._point_in_rect(x, y, marker.icon_rect):
                return marker
            if not icon_only and marker.is_expanded and self._point_in_rect(x, y, marker.popup_rect):
                return marker
        return None

    def _collapse_all_markers(self) -> bool:
        need_refresh = False
        for marker in self.annotations.get(self.current_page, []):
            if marker.is_expanded:
                marker.collapse()
                need_refresh = True
        return need_refresh

    def _style_kind_label(self, marker: AnnotationMarker) -> str:
        from src.models.annotation_preset import _PRESET_BY_ID  # noqa: SLF001

        kid = getattr(marker, "style_kind", "") or ""
        if kid in _PRESET_BY_ID:
            return _PRESET_BY_ID[kid].label
        if getattr(marker, "display_mode", "marker") == "inline":
            return "原位"
        return "标记"

    def _parse_font_size(self, value: str = "") -> int:
        raw = (value or self.font_size_var.get() or "12").strip()
        try:
            size = int(float(raw))
        except ValueError:
            size = 12
        return max(6, min(96, size))

    def _marker_font_family(self, marker: AnnotationMarker) -> str:
        return (
            getattr(marker, "font_family", "")
            or self.settings.annotation.style.font_family
            or "Microsoft YaHei"
        )

    def _display_text_for_marker(self, marker: AnnotationMarker) -> str:
        from src.utils.text_orientation import format_text_for_orientation

        return format_text_for_orientation(
            marker.text,
            getattr(marker, "text_orientation", "horizontal"),
        )

    def _on_text_orientation_change(self, _label: str) -> None:
        self._apply_marker_style_live()

    def _orientation_var_to_value(self) -> str:
        from src.utils.text_orientation import LABEL_TO_ORIENTATION, normalize_orientation

        label = self.text_orientation_var.get() if hasattr(self, "text_orientation_var") else ""
        return normalize_orientation(LABEL_TO_ORIENTATION.get(label, label))

    def _set_orientation_var(self, orientation: str) -> None:
        from src.utils.text_orientation import ORIENTATION_LABELS, normalize_orientation

        if hasattr(self, "text_orientation_var"):
            self.text_orientation_var.set(
                ORIENTATION_LABELS.get(normalize_orientation(orientation), "横向")
            )

    def _sync_annotation_input_style(self) -> None:
        if not hasattr(self, "annotation_input"):
            return
        from src.utils.block_font_size import UI_ANNOTATION_FONT_PT

        family = self.font_family_var.get() if hasattr(self, "font_family_var") else "Microsoft YaHei"
        color = self.color_var.get() if hasattr(self, "color_var") else UITheme.POPUP_TEXT
        try:
            self.annotation_input.configure(
                font=(family, UI_ANNOTATION_FONT_PT),
                text_color=color,
            )
        except tk.TclError:
            pass

    def _apply_marker_style_live(self) -> None:
        if not self.selected_marker:
            return
        self._record_undo_before()
        self.selected_marker.font_size = self._parse_font_size()
        self.selected_marker.font_family = self.font_family_var.get()
        self.selected_marker.text_orientation = self._orientation_var_to_value()
        self._sync_annotation_input_style()
        if self._inline_edit_marker is self.selected_marker and self._inline_edit_widget:
            family, px = self._inline_editor_font(self.selected_marker)
            try:
                self._inline_edit_widget.configure(
                    font=(family, px),
                    text_color=self.selected_marker.color,
                )
            except tk.TclError:
                pass
        else:
            self._show_annotations()
        self.sync_web_preview()

    def _select_marker(self, marker: AnnotationMarker) -> None:
        """选择批注（侧边栏同步）"""
        self.selected_marker = marker
        self.annotation_input.delete("1.0", "end")
        self.annotation_input.insert("1.0", marker.text)
        self.color_var.set(marker.color)
        from src.utils.block_font_size import GENERATED_INLINE_FONT_PT

        fs = getattr(marker, "font_size", 12) or 12
        if (getattr(marker, "original_text", "") or "").strip():
            fs = GENERATED_INLINE_FONT_PT
        self.font_size_var.set(str(fs))
        ff = self._marker_font_family(marker)
        if ff not in self._system_fonts:
            self.font_family_menu.configure(values=[ff] + self._system_fonts)
        self.font_family_var.set(ff)
        self._set_orientation_var(getattr(marker, "text_orientation", "horizontal"))
        self._sync_annotation_input_style()
        self._update_annotation_list()

    def _inline_text_anchor(self, marker: AnnotationMarker) -> str:
        from src.utils.text_orientation import ORIENTATION_VERTICAL, normalize_orientation

        placement = getattr(marker, "placement", "right") or "right"
        vertical = normalize_orientation(getattr(marker, "text_orientation", "")) == ORIENTATION_VERTICAL
        if vertical:
            if placement == "above":
                return "s"
            if placement == "right":
                return "n"
            return "n"
        if placement == "above":
            return "sw"
        if placement == "right":
            return "w"
        return "nw"

    def _inline_editor_font(self, marker: AnnotationMarker) -> tuple:
        family = self._marker_font_family(marker)
        return family, self._inline_font_pixels(marker)

    def _commit_inline_text_edit(self, *, silent: bool = False) -> None:
        """结束画布原位编辑并写回批注。"""
        marker = self._inline_edit_marker
        if not marker or not self._inline_edit_widget:
            self._destroy_inline_editor()
            return

        try:
            text = self._inline_edit_widget.get("1.0", "end-1c").strip()
        except tk.TclError:
            text = marker.text

        if text:
            marker.text = text
            self.annotation_input.delete("1.0", "end")
            self.annotation_input.insert("1.0", marker.text)

        self._destroy_inline_editor()
        self._show_annotations()
        self._update_annotation_list()
        if not silent:
            self._persist_current_file_annotations()
            self.sync_web_preview()
            self.update_status("译文已更新")

    def _cancel_inline_text_edit(self) -> None:
        self._destroy_inline_editor()
        self._show_annotations()

    def _destroy_inline_editor(self) -> None:
        self._inline_edit_opening = False
        if self._inline_edit_toplevel:
            try:
                self._inline_edit_toplevel.destroy()
            except tk.TclError:
                pass
        self._inline_edit_toplevel = None
        self._inline_edit_widget = None
        self._inline_edit_host = None
        self._inline_edit_marker = None

    def _canvas_xy_to_screen(self, cx: float, cy: float) -> Tuple[int, int]:
        self.canvas.update_idletasks()
        ox = self.canvas.winfo_rootx()
        oy = self.canvas.winfo_rooty()
        return (
            int(ox + cx - self.canvas.canvasx(0)),
            int(oy + cy - self.canvas.canvasy(0)),
        )

    def _pick_annotation_color(self) -> None:
        """调色盘选择文字颜色（PPT 式）。"""
        from tkinter import colorchooser

        initial = self.color_var.get() or UITheme.ANNOTATION_COLOR_DEFAULT
        result = colorchooser.askcolor(color=initial, parent=self, title="选择文字颜色")
        if result and result[1]:
            self._set_color(result[1])

    def _deselect_marker(self, *, commit_inline: bool = False) -> None:
        """取消选中（点击画布空白处）。"""
        if self._inline_edit_marker and commit_inline:
            self._commit_inline_text_edit(silent=True)
        elif self._inline_edit_marker:
            self._cancel_inline_text_edit()
        if not self.selected_marker:
            return
        self.selected_marker = None
        self._show_annotations()
        self._update_annotation_list()

    def _place_inline_editor_toplevel(
        self,
        top: ctk.CTkToplevel,
        host: ctk.CTkFrame,
        cx: float,
        cy: float,
        anchor: str,
        width: int,
        height: int,
    ) -> None:
        tw = max(int(width), 120)
        th = max(int(height), 52)
        sx, sy = self._canvas_xy_to_screen(cx, cy)
        if anchor == "sw":
            sx = max(0, sx)
            sy = max(0, sy - th)
        elif anchor == "w":
            sy = max(0, sy - th // 2)
        top.geometry(f"{tw}x{th}+{sx}+{sy}")
        top.deiconify()
        top.lift()
        top.focus_force()

    def _begin_inline_text_edit(self, marker: AnnotationMarker) -> None:
        """双击原位译文：在页面上直接编辑（类似 PPT 文本框）。"""
        if not self.canvas:
            return

        self._commit_inline_text_edit(silent=True)
        self._select_marker(marker)
        self._record_undo_before()
        self._inline_edit_marker = marker

        cx, cy = self._marker_to_canvas(marker)
        anchor = self._inline_text_anchor(marker)
        family, font_px = self._inline_editor_font(marker)

        pad = 8
        if marker.icon_rect and marker.icon_rect[2] > marker.icon_rect[0]:
            w = int(marker.icon_rect[2] - marker.icon_rect[0] + pad * 2)
            h = int(marker.icon_rect[3] - marker.icon_rect[1] + pad * 2 + 36)
        else:
            w = max(160, min(420, len(marker.text) * font_px // 2 + 40))
            h = max(56, font_px + 48)

        w = max(120, min(520, w))
        h = max(52, min(280, h))

        top = ctk.CTkToplevel(self)
        top.withdraw()
        top.transient(self)
        top.attributes("-topmost", True)
        top.overrideredirect(True)
        top.configure(fg_color="#FFFFFF")
        self._inline_edit_toplevel = top
        self._inline_edit_opening = True

        host = ctk.CTkFrame(
            top,
            fg_color="#FFFFFF",
            border_color=marker.color or UITheme.PURPLE_700,
            border_width=2,
            corner_radius=6,
            width=w,
            height=h,
        )
        host.pack_propagate(False)
        host.pack(fill="both", expand=True)

        tool = ctk.CTkFrame(host, fg_color=UITheme.PURPLE_50, height=30, corner_radius=0)
        tool.pack(fill="x")
        tool.pack_propagate(False)

        mini_sizes = ["8", "10", "12", "14", "16", "18", "20", "24"]
        inline_size_menu = ctk.CTkOptionMenu(
            tool,
            values=mini_sizes,
            width=56,
            height=24,
            command=self._on_font_size_preset,
        )
        inline_size_menu.pack(side="left", padx=4, pady=3)
        UITheme.style_option_menu(inline_size_menu)
        inline_font_combo = ctk.CTkComboBox(
            tool,
            values=self._system_fonts[:40],
            width=100,
            variable=self.font_family_var,
            command=self._on_font_family_change,
        )
        inline_font_combo.pack(side="left", padx=2, pady=3)
        UITheme.style_combo(inline_font_combo)

        for color in UITheme.ANNOTATION_COLORS:
            ctk.CTkButton(
                tool,
                text="",
                width=18,
                height=18,
                corner_radius=9,
                command=lambda c=color: self._set_color(c),
                **UITheme.annotation_color_swatch_params(color),
            ).pack(side="left", padx=2, pady=5)

        ctk.CTkButton(
            tool,
            text="🎨",
            width=32,
            height=24,
            command=self._pick_annotation_color,
        ).pack(side="left", padx=2, pady=3)

        ctk.CTkButton(
            tool,
            text="完成",
            width=48,
            height=24,
            command=lambda: self._commit_inline_text_edit(),
        ).pack(side="right", padx=6, pady=3)

        body_h = max(28, h - 34)
        editor = ctk.CTkTextbox(
            host,
            height=body_h,
            wrap="word",
            font=(family, font_px),
            text_color=marker.color or UITheme.PURPLE_800,
            border_width=0,
            fg_color="#FFFFFF",
        )
        editor.pack(fill="both", expand=True, padx=4, pady=(2, 4))
        editor.insert("1.0", marker.text)
        editor.focus_set()
        editor.bind("<Escape>", lambda _e: self._cancel_inline_text_edit())
        editor.bind("<Control-Return>", lambda _e: self._commit_inline_text_edit())

        def on_return_commit(e):
            if e.state & 0x1:
                return
            self._commit_inline_text_edit()
            return "break"

        editor.bind("<Return>", on_return_commit)

        self._inline_edit_host = host
        self._inline_edit_widget = editor
        self._inline_editor_width = w
        self._inline_editor_height = h

        self._show_annotations()
        self._place_inline_editor_toplevel(top, host, cx, cy, anchor, w, h)
        editor.focus_set()
        self._inline_edit_opening = False
        self.update_status("编辑中：Enter 完成 · Esc 取消 · 🎨 调色盘 · 可改字号与颜色")

    def _on_font_size_preset(self, choice: str) -> None:
        self.font_size_var.set(choice)
        self._on_font_size_apply()

    def _on_font_size_apply(self) -> None:
        self._apply_marker_style_live()

    def _on_font_family_change(self, choice: str) -> None:
        if choice:
            self.font_family_var.set(choice)
        self._apply_marker_style_live()

    def _expand_marker(self, marker: AnnotationMarker) -> None:
        for m in self.annotations.get(self.current_page, []):
            if m != marker:
                m.collapse()
        marker.toggle_expand()
        self._select_marker(marker)
        self._show_annotations()

    def _on_canvas_press(self, event) -> None:
        """鼠标按下：选中批注并准备拖动"""
        if not self.canvas:
            return

        x = self.canvas.canvasx(event.x)
        y = self.canvas.canvasy(event.y)
        self._press_pos = (x, y)
        self._did_drag = False
        self._drag_undo_snap = None

        if self.canvas_tool == "eraser":
            self._record_undo_before()
            self._erase_ink_at(x, y)
            self._did_drag = True
            return

        if self.canvas_tool in ("pen", "highlighter"):
            self._record_undo_before()
            self._start_ink_stroke(x, y)
            self._did_drag = True
            return

        if (
            self._inline_edit_marker
            and self._inline_edit_host
            and not self._inline_edit_opening
        ):
            try:
                wx = self.canvas.winfo_rootx() + int(event.x)
                wy = self.canvas.winfo_rooty() + int(event.y)
                hx = self._inline_edit_host.winfo_rootx()
                hy = self._inline_edit_host.winfo_rooty()
                hw = max(self._inline_edit_host.winfo_width(), 1)
                hh = max(self._inline_edit_host.winfo_height(), 1)
                if not (hx <= wx <= hx + hw and hy <= wy <= hy + hh):
                    self._commit_inline_text_edit(silent=True)
                    self.selected_marker = None
                    self._update_annotation_list()
                    return
            except tk.TclError:
                self._commit_inline_text_edit(silent=True)
                self.selected_marker = None
                self._update_annotation_list()
                return

        clicked_marker = self._find_marker_at(x, y, icon_only=True)
        if clicked_marker:
            self._drag_undo_snap = self._capture_undo_snapshot()
            self.selected_marker = clicked_marker
            pdf_x, pdf_y = self._canvas_to_pdf(x, y)
            clicked_marker.drag_data = {
                "x": pdf_x - clicked_marker.x,
                "y": pdf_y - clicked_marker.y,
            }
            self._select_marker(clicked_marker)
            self.dragging = True
            if getattr(clicked_marker, "display_mode", "marker") != "inline":
                self._show_annotations()
        else:
            if self.selected_marker and not self._inline_edit_marker:
                self._deselect_marker()
            else:
                self.selected_marker = None
            self.dragging = False

    def _on_canvas_drag(self, event) -> None:
        """鼠标拖动：移动批注"""
        if self._current_ink_stroke:
            x = self.canvas.canvasx(event.x)
            y = self.canvas.canvasy(event.y)
            self._extend_ink_stroke(x, y)
            return
        if self.canvas_tool == "eraser":
            x = self.canvas.canvasx(event.x)
            y = self.canvas.canvasy(event.y)
            self._erase_ink_at(x, y)
            return
        if self._inline_edit_marker:
            return
        if not self.dragging or not self.selected_marker:
            return

        x = self.canvas.canvasx(event.x)
        y = self.canvas.canvasy(event.y)

        if self._press_pos:
            dx = abs(x - self._press_pos[0])
            dy = abs(y - self._press_pos[1])
            if dx > 3 or dy > 3:
                self._did_drag = True

        pdf_x, pdf_y = self._canvas_to_pdf(x, y)
        marker = self.selected_marker
        page = self.pdf_doc[self.current_page]
        nx = pdf_x - marker.drag_data["x"]
        ny = pdf_y - marker.drag_data["y"]
        if getattr(marker, "display_mode", "marker") == "inline":
            marker.x = int(max(0, min(nx, page.rect.width - 1)))
            marker.y = int(max(0, min(ny, page.rect.height - 1)))
        else:
            size = AnnotationMarker.MARKER_SIZE
            marker.x = int(max(0, min(nx, page.rect.width - size)))
            marker.y = int(max(0, min(ny, page.rect.height - size)))

        self._show_annotations()

    def _on_canvas_release(self, event) -> None:
        """鼠标释放：区分点击与拖动"""
        if not self.canvas:
            return

        x = self.canvas.canvasx(event.x)
        y = self.canvas.canvasy(event.y)

        if self._current_ink_stroke:
            self._extend_ink_stroke(x, y)
            self._finish_ink_stroke()
            self._press_pos = None
            self._did_drag = False
            return

        if self._did_drag and self.canvas_tool == "eraser":
            self.dragging = False
            self._press_pos = None
            self._did_drag = False
            return

        if self._did_drag:
            self.dragging = False
            self._flush_drag_undo_snap()
            self._persist_current_file_annotations()
            msg = (
                "译文位置已更新"
                if getattr(self.selected_marker, "display_mode", "") == "inline"
                else "批注位置已更新"
            )
            self.update_status(msg)
            self.sync_web_preview()
            return

        clicked_marker = self._find_marker_at(x, y)

        if clicked_marker:
            if getattr(clicked_marker, "display_mode", "marker") == "inline":
                self._select_marker(clicked_marker)
                if self._inline_edit_marker is not clicked_marker:
                    self._show_annotations()
            elif clicked_marker.is_expanded and clicked_marker.close_rect and self._point_in_rect(
                x, y, clicked_marker.close_rect
            ):
                clicked_marker.collapse()
                self.selected_marker = None
                self._show_annotations()
                self._update_annotation_list()
            else:
                self._select_marker(clicked_marker)
        else:
            if self._inline_edit_marker:
                pass
            elif self.selected_marker:
                self._deselect_marker()
            elif self._collapse_all_markers():
                self.selected_marker = None
                self._show_annotations()
                self._update_annotation_list()

            if hasattr(self, "_adding_annotation") and self._adding_annotation:
                self._create_annotation_at(x, y)
                self._adding_annotation = False
                self.mode_hint.configure(text="拖动标记可移动，双击查看批注，点击空白处关闭")

        self.dragging = False
        self._press_pos = None
        self._did_drag = False

    def _on_canvas_double_click(self, event) -> None:
        """双击：原位译文进入编辑；普通标记展开弹层"""
        if not self.canvas:
            return

        x = self.canvas.canvasx(event.x)
        y = self.canvas.canvasy(event.y)
        marker = self._find_marker_at(x, y, icon_only=True)
        if not marker:
            return

        if getattr(marker, "display_mode", "marker") == "inline":
            self.dragging = False
            self._did_drag = False
            edit_target = marker
            self.after(50, lambda m=edit_target: self._begin_inline_text_edit(m))
            return

        if marker.is_expanded:
            marker.collapse()
            self.selected_marker = None
            self._show_annotations()
            self._update_annotation_list()
        else:
            self._expand_marker(marker)

    def _get_add_preset(self):
        from src.models.annotation_preset import get_preset_by_label

        label = self.add_preset_var.get() if hasattr(self, "add_preset_var") else "原位译文"
        preset = self._pending_add_preset or get_preset_by_label(label)
        if preset.id == "custom":
            return preset, {
                "display_mode": "inline",
                "color": self.color_var.get(),
                "font_size": self._parse_font_size(),
                "font_family": self.font_family_var.get(),
                "text_orientation": self._orientation_var_to_value(),
                "placement": "right",
                "style_kind": "custom",
            }
        return preset, {
            "display_mode": preset.display_mode,
            "color": preset.color,
            "font_size": preset.font_size,
            "font_family": preset.font_family,
            "text_orientation": self._orientation_var_to_value(),
            "placement": preset.placement,
            "style_kind": preset.id,
        }

    def _apply_preset_to_sidebar(self, preset) -> None:
        """将样式预设同步到批注内容区（「当前样式」除外）。"""
        if preset.id == "custom":
            return
        if not hasattr(self, "_system_fonts"):
            from src.utils.system_fonts import get_system_font_families

            self._system_fonts = get_system_font_families()
        self.color_var.set(preset.color)
        self.font_size_var.set(str(preset.font_size))
        if preset.font_family in self._system_fonts:
            self.font_family_var.set(preset.font_family)
        else:
            self.font_family_menu.configure(values=[preset.font_family] + self._system_fonts)
            self.font_family_var.set(preset.font_family)
        self._sync_annotation_input_style()

    def _on_add_preset_change(self, label: str) -> None:
        from src.models.annotation_preset import get_preset_by_label

        preset = get_preset_by_label(label)
        self._pending_add_preset = preset
        self._apply_preset_to_sidebar(preset)
        if hasattr(self, "mode_hint") and not getattr(self, "_adding_annotation", False):
            self.mode_hint.configure(text=f"添加类型：{preset.label} — {preset.hint}")

    def _clear_annotation_editor(self) -> None:
        """清空批注编辑区并取消选中（用于新建批注，避免沿用 AI/旧批注正文）"""
        self._cancel_inline_text_edit()
        self.selected_marker = None
        self.annotation_input.delete("1.0", "end")
        self._update_annotation_list()

    def _add_annotation_mode(self) -> None:
        """选择样式后进入添加批注模式"""
        preset, _style = self._get_add_preset()
        self._clear_annotation_editor()
        self._apply_preset_to_sidebar(preset)
        self._adding_annotation = True
        self.mode_hint.configure(
            text=f"请在页面点击放置「{preset.label}」"
        )

    def _create_annotation_at(self, x: int, y: int) -> None:
        """在指定位置创建批注（按所选样式类型）"""
        preset, style = self._get_add_preset()
        text = self.annotation_input.get("1.0", "end-1c").strip()
        if not text:
            text = "新批注" if preset.id == "marker" else "请输入译文"

        pdf_x, pdf_y = self._canvas_to_pdf(x, y)

        marker = AnnotationMarker(
            int(pdf_x),
            int(pdf_y),
            text,
            style["color"],
            display_mode=style["display_mode"],
            placement=style["placement"],
            font_size=style["font_size"],
            font_family=style["font_family"],
            text_orientation=style.get("text_orientation", "horizontal"),
            style_kind=style["style_kind"],
        )

        if self.current_page not in self.annotations:
            self.annotations[self.current_page] = []

        self._record_undo_before()
        self.annotations[self.current_page].append(marker)
        self._select_marker(marker)
        self._show_annotations()

        self.update_status(f"已添加「{preset.label}」批注")
        self.sync_web_preview()

        if marker.display_mode == "inline":
            self.after(80, lambda m=marker: self._begin_inline_text_edit(m))

    def _save_annotation(self) -> None:
        """保存批注内容"""
        if not self.selected_marker:
            show_warning(self, "警告", "请先选择一个批注")
            return

        self._record_undo_before()
        text = self.annotation_input.get("1.0", "end-1c")
        if not text.strip():
            show_warning(self, "警告", "批注内容不能为空")
            return

        self.selected_marker.text = text
        self.selected_marker.color = self.color_var.get()
        self.selected_marker.font_size = self._parse_font_size()
        self.selected_marker.font_family = self.font_family_var.get()
        self.selected_marker.text_orientation = self._orientation_var_to_value()
        self.selected_marker.refresh_expanded_size()

        self._show_annotations()
        self._update_annotation_list()

        self.update_status("批注已保存")
        self.sync_web_preview()

    def _remove_marker(self, marker: AnnotationMarker) -> None:
        """从当前页移除一条批注并刷新界面。"""
        self._record_undo_before()
        if self._inline_edit_marker is marker:
            self._cancel_inline_text_edit()

        markers = self.annotations.get(self.current_page, [])
        if marker in markers:
            markers.remove(marker)

        if self.selected_marker is marker:
            self.selected_marker = None
            self.annotation_input.delete("1.0", "end")

        self._show_annotations()
        self._update_annotation_list()
        self._persist_current_file_annotations()
        self.sync_web_preview()

    def _delete_marker_item(self, marker: AnnotationMarker) -> None:
        """列表行上的单条删除。"""
        self._remove_marker(marker)
        self.update_status("批注已删除")

    def _delete_all_page_annotations(self) -> None:
        """删除当前页全部批注。"""
        markers = self.annotations.get(self.current_page, [])
        if not markers:
            show_warning(self, "提示", "当前页没有批注")
            return

        if not ask_yes_no(
            self,
            "删除全部批注",
            f"确定删除第 {self.current_page + 1} 页的全部 {len(markers)} 条批注吗？\n删除后可用 Ctrl+Z 撤销。",
            width=420,
        ):
            return

        self._cancel_inline_text_edit()
        self._record_undo_before()
        self.annotations[self.current_page] = []
        self.selected_marker = None
        self.annotation_input.delete("1.0", "end")
        self._show_annotations()
        self._update_annotation_list()
        self._persist_current_file_annotations()
        self.sync_web_preview()
        self.update_status(f"已删除第 {self.current_page + 1} 页全部批注")

    def _is_text_entry_focused(self) -> bool:
        """焦点在可编辑文本框内时不应把退格当作删除批注。"""
        w = self.focus_get()
        while w is not None:
            try:
                cls = w.winfo_class()
            except tk.TclError:
                break
            if cls in ("Entry", "Text", "TEntry"):
                return True
            w = getattr(w, "master", None)
        return False

    def _on_backspace_delete_annotation(self, event=None) -> Optional[str]:
        """选中批注时按 Backspace/Delete 直接删除。"""
        if self._is_text_entry_focused():
            return None
        if not self.selected_marker:
            return None
        self._remove_marker(self.selected_marker)
        self.update_status("批注已删除")
        return "break"

    def _delete_annotation(self) -> None:
        """删除当前选中的批注（批注内容区按钮）。"""
        if not self.selected_marker:
            show_warning(self, "警告", "请先选择一个批注")
            return
        self._remove_marker(self.selected_marker)
        self.update_status("批注已删除")

    def _set_color(self, color: str) -> None:
        """设置批注文字颜色（编辑中不重算窗口大小，避免点色块框变大）。"""
        self.color_var.set(color)
        if not self.selected_marker:
            return
        self.selected_marker.color = color
        if self._inline_edit_marker is self.selected_marker and self._inline_edit_widget:
            try:
                self._inline_edit_widget.configure(text_color=color)
                if self._inline_edit_host:
                    self._inline_edit_host.configure(border_color=color)
            except tk.TclError:
                pass
            self.sync_web_preview()
            return
        self._sync_annotation_input_style()
        self._show_annotations()
        self.sync_web_preview()

    def _update_annotation_list(self) -> None:
        """更新批注列表"""
        # 清空列表
        for widget in self.annotation_list.winfo_children():
            widget.destroy()

        markers = self.annotations.get(self.current_page, [])

        list_pad = UITheme.PAD_SM
        for idx, marker in enumerate(markers):
            frame = ctk.CTkFrame(self.annotation_list, fg_color="transparent", cursor="hand2")
            frame.pack(fill="x", pady=3, padx=0)
            frame.bind("<Button-1>", lambda _e, m=marker: self._select_marker(m))
            UITheme.style_annotation_row(frame, selected=(marker == self.selected_marker))

            inner = ctk.CTkFrame(frame, fg_color="transparent", cursor="hand2")
            inner.pack(fill="x", padx=list_pad, pady=6)
            inner.bind("<Button-1>", lambda _e, m=marker: self._select_marker(m))
            inner.grid_columnconfigure(1, weight=1)

            del_btn = ctk.CTkButton(
                inner,
                text="✕",
                width=30,
                height=28,
                corner_radius=UITheme.RADIUS_SM,
                command=lambda m=marker: self._delete_marker_item(m),
            )
            del_btn.grid(row=0, column=2, padx=(6, 0), sticky="e")
            UITheme.style_icon_dismiss(del_btn)

            color_indicator = ctk.CTkLabel(
                inner,
                text="●",
                text_color=marker.color,
                font=("Arial", 16),
                width=20,
            )
            color_indicator.grid(row=0, column=0, padx=(0, 6), sticky="w")
            color_indicator.bind("<Button-1>", lambda _e, m=marker: self._select_marker(m))

            kind_label = self._style_kind_label(marker)
            preview_text = format_annotation_list_preview(marker.text)
            if kind_label:
                preview_text = f"[{kind_label}] {preview_text}"
            text_label = ctk.CTkLabel(
                inner,
                text=preview_text,
                anchor="w",
                justify="left",
                font=UITheme.font_body(),
                text_color=UITheme.TEXT,
                cursor="hand2",
            )
            text_label.grid(row=0, column=1, sticky="ew")
            text_label.bind("<Button-1>", lambda _e, m=marker: self._select_marker(m))

    def _prev_page(self) -> None:
        """上一页"""
        if not self.pdf_doc:
            return

        if self.current_page > 0:
            self._collapse_all_markers()
            self.current_page -= 1
            self._render_page()

    def _next_page(self) -> None:
        """下一页"""
        if not self.pdf_doc:
            return

        if self.current_page < self.total_pages - 1:
            self._collapse_all_markers()
            self.current_page += 1
            self._render_page()

    def _zoom_in(self) -> None:
        """放大"""
        if self.zoom_level < 3.0:
            self.zoom_level += 0.1
            self._render_page()
            self.zoom_label.configure(text=f"{int(self.zoom_level * 100)}%")

    def _zoom_out(self) -> None:
        """缩小"""
        if self.zoom_level > 0.5:
            self.zoom_level -= 0.1
            self._render_page()
            self.zoom_label.configure(text=f"{int(self.zoom_level * 100)}%")

    def _zoom_reset(self) -> None:
        """重置缩放"""
        self.zoom_level = 1.0
        self._render_page()
        self.zoom_label.configure(text="100%")

    def _remove_current_file(self) -> None:
        """移除当前选中的文件"""
        if self.current_file_index < 0 or not self.selected_files:
            self.update_status("没有可移除的文件")
            return
        self._remove_file(self.current_file_index)

    def _remove_file(self, index: int) -> None:
        """从列表中移除已导入的文件（不删除磁盘原文件）"""
        if not (0 <= index < len(self.selected_files)):
            return

        removed = self.selected_files[index]
        file_name = os.path.basename(removed)
        key = self._file_key(removed)
        has_annotations = bool(self.annotations_by_file.get(key))

        if has_annotations:
            if not ask_yes_no(
                self,
                "确认移除",
                f"确定从列表中移除「{file_name}」吗？\n该文件的批注也会一并删除（不会删除磁盘上的原文件）。",
            ):
                return

        self._persist_current_file_annotations()
        self.selected_files.pop(index)
        self.annotations_by_file.pop(key, None)
        self.converted_pdf_paths.pop(key, None)

        if index == self.current_file_index:
            if self.pdf_doc:
                self.pdf_doc.close()
                self.pdf_doc = None

            if self.selected_files:
                new_index = min(index, len(self.selected_files) - 1)
                self._on_file_select(new_index)
            else:
                self.current_file_index = -1
                self.total_pages = 0
                self.current_page = 0
                self.annotations = {}
                self.selected_marker = None

                self._destroy_preview_shell()
                self.file_hint.configure(text="导入文件")

        elif index < self.current_file_index:
            self.current_file_index -= 1

        self.update_file_list(self.selected_files)
        self.update_status(f"已移除: {file_name}")
        self.sync_web_preview()
        self._flush_persist()
