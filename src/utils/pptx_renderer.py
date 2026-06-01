"""不依赖 Office 的 PPTX 幻灯片渲染（python-pptx + Pillow）"""
from __future__ import annotations

import base64
import io
import os
import textwrap
from typing import Optional

import fitz
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

RENDER_VERSION = 2
SLIDE_WIDTH_PX = 1920


def _load_font(size: int = 16) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
        r"C:\Windows\Fonts\arial.ttf",
        "/System/Library/Fonts/PingFang.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.isfile(path):
            try:
                return ImageFont.truetype(path, size=size)
            except OSError:
                continue
    return ImageFont.load_default()


def _scale(value: int, src: int, dst: int) -> int:
    if src <= 0:
        return 0
    return int(value * dst / src)


def _slide_size_pt(prs: Presentation) -> tuple[float, float]:
    width_pt = prs.slide_width / 914400 * 72
    height_pt = prs.slide_height / 914400 * 72
    return width_pt, height_pt


def _shape_box(
    shape,
    slide_w: int,
    slide_h: int,
    px_w: int,
    px_h: int,
    offset_x: int = 0,
    offset_y: int = 0,
) -> tuple[int, int, int, int]:
    left = shape.left + offset_x
    top = shape.top + offset_y
    x1 = _scale(left, slide_w, px_w)
    y1 = _scale(top, slide_h, px_h)
    x2 = _scale(left + shape.width, slide_w, px_w)
    y2 = _scale(top + shape.height, slide_h, px_h)
    x1 = max(0, min(x1, px_w - 1))
    y1 = max(0, min(y1, px_h - 1))
    x2 = max(x1 + 1, min(x2, px_w))
    y2 = max(y1 + 1, min(y2, px_h))
    return x1, y1, x2, y2


def _estimate_font_size(shape, box_h: int) -> int:
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        return max(12, min(48, int(run.font.size.pt)))
    except Exception:
        pass
    return max(14, min(36, box_h // 5))


def _text_frame_box(
    shape,
    box: tuple[int, int, int, int],
    slide_w: int,
    slide_h: int,
    px_w: int,
    px_h: int,
) -> tuple[int, int, int, int]:
    x1, y1, x2, y2 = box
    try:
        tf = shape.text_frame
        ml = _scale(int(tf.margin_left or 0), slide_w, px_w)
        mt = _scale(int(tf.margin_top or 0), slide_h, px_h)
        mr = _scale(int(tf.margin_right or 0), slide_w, px_w)
        mb = _scale(int(tf.margin_bottom or 0), slide_h, px_h)
        return (
            min(x1 + ml, x2 - 1),
            min(y1 + mt, y2 - 1),
            max(x2 - mr, x1 + 1),
            max(y2 - mb, y1 + 1),
        )
    except Exception:
        return x1 + 4, y1 + 4, x2 - 4, y2 - 4


def _draw_text_in_box(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: tuple[int, int, int, int],
    font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    fill: str = "#222222",
) -> None:
    x1, y1, x2, y2 = box
    width = max(x2 - x1, 40)
    lines: list[str] = []
    avg_char_w = max(getattr(font, "size", 16) * 0.55, 6)
    wrap_width = max(8, int(width / avg_char_w))

    for paragraph in text.splitlines():
        paragraph = paragraph.strip()
        if not paragraph:
            lines.append("")
            continue
        lines.extend(textwrap.wrap(paragraph, width=wrap_width) or [""])

    line_height = getattr(font, "size", 16) + 6
    y = y1
    max_y = y2
    for line in lines:
        if y + line_height > max_y:
            if y <= y1:
                draw.text((x1, y), line[: wrap_width - 1] + "…", fill=fill, font=font)
            break
        draw.text((x1, y), line, fill=fill, font=font)
        y += line_height


def _render_shape(
    shape,
    canvas: Image.Image,
    draw: ImageDraw.ImageDraw,
    slide_w: int,
    slide_h: int,
    px_w: int,
    px_h: int,
    offset_x: int = 0,
    offset_y: int = 0,
) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        gx = offset_x + shape.left
        gy = offset_y + shape.top
        for child in shape.shapes:
            _render_shape(child, canvas, draw, slide_w, slide_h, px_w, px_h, gx, gy)
        return

    box = _shape_box(shape, slide_w, slide_h, px_w, px_h, offset_x, offset_y)
    x1, y1, x2, y2 = box

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
            w, h = max(x2 - x1, 1), max(y2 - y1, 1)
            pic = pic.resize((w, h), Image.Resampling.LANCZOS)
            canvas.paste(pic, (x1, y1), pic)
        except Exception:
            draw.rectangle(box, outline="#cccccc", fill="#f5f5f5")
        return

    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            inner = _text_frame_box(shape, box, slide_w, slide_h, px_w, px_h)
            font_size = _estimate_font_size(shape, inner[3] - inner[1])
            font = _load_font(font_size)
            _draw_text_in_box(draw, text, inner, font)
        return

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            if rows and cols:
                cell_h = max((y2 - y1) // rows, 1)
                cell_w = max((x2 - x1) // cols, 1)
                font = _load_font(max(11, min(18, cell_h // 2)))
                for r, row in enumerate(table.rows):
                    for c, cell in enumerate(row.cells):
                        cx1 = x1 + c * cell_w
                        cy1 = y1 + r * cell_h
                        draw.rectangle(
                            (cx1, cy1, cx1 + cell_w, cy1 + cell_h),
                            outline="#dddddd",
                        )
                        cell_text = (cell.text or "").strip()
                        if cell_text:
                            _draw_text_in_box(
                                draw,
                                cell_text,
                                (cx1 + 2, cy1 + 2, cx1 + cell_w - 2, cy1 + cell_h - 2),
                                font,
                            )
        except Exception:
            pass


def _render_slide(slide, slide_w: int, slide_h: int) -> Image.Image:
    px_w = SLIDE_WIDTH_PX
    px_h = _scale(slide_h, slide_w, px_w)
    image = Image.new("RGB", (px_w, px_h), "white")
    draw = ImageDraw.Draw(image)

    for shape in slide.shapes:
        _render_shape(shape, image, draw, slide_w, slide_h, px_w, px_h)

    return image


def _image_to_pdf_page(doc: fitz.Document, image: Image.Image, width_pt: float, height_pt: float) -> None:
    """将幻灯片图片写入 PDF 页（使用标准 point 尺寸，避免预览裁切）"""
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    page = doc.new_page(width=width_pt, height=height_pt)
    page.insert_image(fitz.Rect(0, 0, width_pt, height_pt), stream=buf.getvalue())


def render_pptx_to_pdf(source: str, output_pdf: str) -> None:
    """将 PPTX 渲染为 PDF（无需 Office）"""
    prs = Presentation(source)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    if not prs.slides:
        raise ValueError("PPT 中没有幻灯片")

    width_pt, height_pt = _slide_size_pt(prs)
    doc = fitz.open()
    try:
        for slide in prs.slides:
            image = _render_slide(slide, slide_w, slide_h)
            _image_to_pdf_page(doc, image, width_pt, height_pt)
        doc.save(output_pdf)
    finally:
        doc.close()


def can_render_pptx(source: str) -> bool:
    return source.lower().endswith(".pptx")


def extract_pptx_slide_text(source: str, slide_index: int) -> str:
    """从 PPTX 单页提取文字（标题、正文、备注）"""
    prs = Presentation(source)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return ""

    slide = prs.slides[slide_index]
    texts: list[str] = []

    def _collect_shapes(shapes, prefix=""):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _collect_shapes(shape.shapes, prefix)
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)

    _collect_shapes(slide.shapes)

    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame and notes_frame.text.strip():
            texts.append(f"[备注] {notes_frame.text.strip()}")

    return "\n".join(texts)


def render_pptx_slide_to_image(
    source: str,
    slide_index: int,
) -> tuple[bytes, str, float, float]:
    """将 PPTX 单页幻灯片渲染为 PNG"""
    prs = Presentation(source)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"幻灯片索引超出范围: {slide_index} ({len(prs.slides)} 页)")

    slide = prs.slides[slide_index]
    image = _render_slide(slide, prs.slide_width, prs.slide_height)
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    png_bytes = buf.getvalue()
    b64_str = base64.b64encode(png_bytes).decode("utf-8")
    width_pt, height_pt = _slide_size_pt(prs)
    return png_bytes, b64_str, width_pt, height_pt
