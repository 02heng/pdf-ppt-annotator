from pptx import Presentation
import os
from src.models.document import Document
from src.models.page import Page
from .base_parser import BaseParser


class PPTParser(BaseParser):
    """PPT 文档解析器"""

    def parse(self, file_path: str) -> Document:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        if not self.supports_format(file_path):
            raise ValueError(f"不支持的文件格式: {file_path}")

        doc = Document(file_path=file_path, file_type="ppt")

        try:
            prs = Presentation(file_path)

            if prs.core_properties:
                doc.title = prs.core_properties.title
                doc.author = prs.core_properties.author

            for slide_num, slide in enumerate(prs.slides):
                text_content = self._extract_slide_text(slide)
                images = self._extract_slide_images(slide)

                width = prs.slide_width / 914400
                height = prs.slide_height / 914400

                page = Page(
                    page_number=slide_num + 1,
                    content=text_content,
                    images=images,
                    width=width,
                    height=height,
                )
                doc.add_page(page)

        except Exception as e:
            raise ValueError(f"PPT 解析错误: {str(e)}")

        return doc

    def supports_format(self, file_path: str) -> bool:
        return file_path.lower().endswith((".ppt", ".pptx"))

    def _extract_slide_text(self, slide) -> str:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)

        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                texts.append(f"[备注] {notes_slide.notes_text_frame.text}")

        return "\n".join(texts)

    def _extract_slide_images(self, slide) -> list:
        images = []
        for shape in slide.shapes:
            if shape.shape_type == 13:
                images.append(f"image_{shape.shape_id}")
        return images
